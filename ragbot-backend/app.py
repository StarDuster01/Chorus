import os
import json
import base64
import uuid
import datetime
import sys  # Add sys import
from datetime import UTC  # Import UTC for timezone-aware datetime objects
import tempfile
import bcrypt
import jwt
import re
from werkzeug.utils import secure_filename
from flask import Flask, request, jsonify, send_from_directory, send_file
from flask_cors import CORS
from dotenv import load_dotenv
import openai
import anthropic
import requests  # For Groq API
import chromadb
from chromadb.utils import embedding_functions
from PyPDF2 import PdfReader
import docx
from pptx import Presentation
from PIL import Image  # Add this import for image processing
from image_processor import ImageProcessor
import faiss
import pytesseract
import fitz  # PyMuPDF
from pdf2image import convert_from_path
import io
import numpy as np

# Load environment variables
load_dotenv()

# Set OpenMP environment variable to avoid library conflicts
os.environ["KMP_DUPLICATE_LIB_OK"] = "TRUE"

# Define supported file extensions
text_extensions = ['.pdf', '.docx', '.txt', '.pptx']
image_extensions = ['.jpg', '.jpeg', '.png', '.gif', '.webp', '.bmp']

# Set up necessary directories
app_base_dir = os.path.dirname(os.path.abspath(__file__))
DATA_FOLDER = os.path.join(app_base_dir, "data")
os.makedirs(DATA_FOLDER, exist_ok=True)

# Configure OCR
try:
    # On Windows, you might need to set the path to Tesseract executable
    if os.name == 'nt':  # If on Windows
        pytesseract.pytesseract.tesseract_cmd = os.getenv('TESSERACT_PATH', r'C:\Program Files\Tesseract-OCR\tesseract.exe')
except Exception as e:
    print(f"Warning: Could not configure pytesseract: {str(e)}")

# Define a function to extract text from images using OCR
def extract_text_from_image(img):
    """Extract text from an image using OCR
    
    Args:
        img: PIL Image object or path to image
        
    Returns:
        str: Extracted text
    """
    try:
        if isinstance(img, str):
            # If img is a path
            img = Image.open(img)
        
        # Convert image to RGB if it's not
        if img.mode != 'RGB':
            img = img.convert('RGB')
            
        # Apply image preprocessing to improve OCR
        # Convert to grayscale
        gray = img.convert('L')
        
        # Apply threshold to create a black and white image
        threshold = 128
        bw = gray.point(lambda x: 0 if x < threshold else 255, '1')
        
        # Extract text using pytesseract
        text = pytesseract.image_to_string(bw)
        
        return text
    except Exception as e:
        print(f"Warning: Failed to extract text from image: {str(e)}")
        return ""

def extract_text_from_pdf(pdf_path):
    """Extract text from a PDF including text in images
    
    Args:
        pdf_path: Path to the PDF file
        
    Returns:
        str: Extracted text
    """
    text = ""
    image_text = ""
    
    # First, extract regular text using PyPDF2
    with open(pdf_path, 'rb') as file:
        reader = PdfReader(file)
        for page in reader.pages:
            text += page.extract_text() + "\n"
    
    # Then, extract text from images using PyMuPDF and OCR
    try:
        doc = fitz.open(pdf_path)
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            
            # Get images from the page
            image_list = page.get_images(full=True)
            
            # Process each image
            for img_index, img in enumerate(image_list):
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                
                # Load image using PIL
                img = Image.open(io.BytesIO(image_bytes))
                
                # Extract text using OCR
                img_text = extract_text_from_image(img)
                if img_text.strip():
                    image_text += f"Image {page_num+1}.{img_index+1} text: {img_text}\n"
        
        doc.close()
    except Exception as e:
        print(f"Warning: Failed to extract images from PDF: {str(e)}")
        
        # Fall back to pdf2image if PyMuPDF fails
        try:
            # Convert PDF to images
            images = convert_from_path(pdf_path)
            
            # Process each page image
            for i, img in enumerate(images):
                # Extract text using OCR
                img_text = extract_text_from_image(img)
                if img_text.strip():
                    image_text += f"Page {i+1} image text: {img_text}\n"
        except Exception as e2:
            print(f"Warning: Failed to use pdf2image fallback: {str(e2)}")
    
    # Combine regular text and image text
    combined_text = text
    if image_text:
        combined_text += "\n\nExtracted from images:\n" + image_text
        
    return combined_text

def extract_text_from_docx(docx_path):
    """Extract text from a Word document
    
    Args:
        docx_path: Path to the docx file
        
    Returns:
        str: Extracted text
    """
    text = ""
    doc = docx.Document(docx_path)
    
    # Extract text from paragraphs
    for para in doc.paragraphs:
        text += para.text + "\n"
        
    # Extract text from tables
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                text += cell.text + " "
            text += "\n"
    
    # Unfortunately, direct image extraction from docx is complex
    # We don't have a straightforward way to extract images for OCR
    # A complete solution would require unpacking the docx and processing embedded images
    
    return text

def extract_text_from_pptx(pptx_path):
    """Extract text from a PowerPoint presentation including text in images
    
    Args:
        pptx_path: Path to the pptx file
        
    Returns:
        str: Extracted text
    """
    summary_lines = []
    image_text = ""
    image_captions = ""
    prs = Presentation(pptx_path)
    
    # Log the total number of slides
    slide_count = len(prs.slides)
    print(f"Processing PowerPoint with {slide_count} slides")
    
    # Add presentation metadata to summary
    summary_lines.append("## PRESENTATION SUMMARY ##")
    summary_lines.append(f"Total Slides: {slide_count}")
    summary_lines.append("")
    
    # Get or create ImageProcessor for captioning
    try:
        img_processor = ImageProcessor(os.path.join(os.path.dirname(os.path.abspath(__file__)), "data"))
    except Exception as e:
        print(f"Warning: Could not initialize ImageProcessor: {str(e)}")
        img_processor = None
    
    # Create temp directory for extracting images
    with tempfile.TemporaryDirectory() as temp_dir:
        # Process each slide
        for slide_num, slide in enumerate(prs.slides):
            # Add clear slide separator with prominent marker
            summary_lines.append(f"## SLIDE {slide_num+1} ##")
            slide_has_content = False
            
            # Try to get slide title
            slide_title = ""
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text and hasattr(shape, 'is_title') and shape.is_title:
                    slide_title = shape.text_frame.text.strip()
                    summary_lines.append(f"TITLE: {slide_title}")
                    slide_has_content = True
                    break
            
            # Add a section for text content
            text_content = []
            
            # Get text from all shapes with text
            for shape in slide.shapes:
                # Get text from text frames
                if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    if shape.text_frame and shape.text_frame.text and shape.text_frame.text.strip():
                        # Skip if we already included this as the title
                        text_content_item = shape.text_frame.text.strip()
                        if text_content_item != slide_title:
                            text_content.append(text_content_item)
                            slide_has_content = True
                # Also check for text attribute (some shapes have this instead)
                elif hasattr(shape, "text") and shape.text and shape.text.strip():
                    # Skip if we already included this as the title
                    text_content_item = shape.text.strip()
                    if text_content_item != slide_title:
                        text_content.append(text_content_item)
                        slide_has_content = True
                
                # Get text from tables if present
                if shape.has_table:
                    table_content = ["TABLE CONTENT:"]
                    has_table_content = False
                    for row in shape.table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text_frame.text:
                                row_text.append(cell.text_frame.text.strip())
                                has_table_content = True
                        if row_text:
                            table_content.append(" | ".join(row_text))
                    
                    if has_table_content:
                        text_content.extend(table_content)
                        slide_has_content = True
            
            # Add text content section if there's any content
            if text_content:
                summary_lines.append("\nCONTENT:")
                for item in text_content:
                    summary_lines.append(item)
            
            # Add a section for image content
            slide_images = []
            
            # Extract image if this shape is a picture
            for shape_index, shape in enumerate(slide.shapes):
                if hasattr(shape, "shape_type") and shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    try:
                        # Save image to temp file
                        image_path = os.path.join(temp_dir, f"slide_{slide_num+1}_img_{uuid.uuid4()}.png")
                        with open(image_path, 'wb') as f:
                            f.write(shape.image.blob)
                        
                        # Process the image in two ways:
                        image_info = {
                            "index": shape_index,
                            "path": image_path,
                            "ocr_text": "",
                            "caption": ""
                        }
                        
                        # 1. Extract text using OCR
                        img_text = extract_text_from_image(image_path)
                        if img_text.strip():
                            image_info["ocr_text"] = img_text.strip()
                            slide_has_content = True
                        
                        # 2. Generate image caption if ImageProcessor is available
                        if img_processor:
                            try:
                                caption = img_processor.generate_caption(image_path)
                                if caption and caption != "No caption available":
                                    image_info["caption"] = caption
                                    slide_has_content = True
                            except Exception as e:
                                print(f"Warning: Failed to generate caption for image on slide {slide_num+1}: {str(e)}")
                        
                        slide_images.append(image_info)
                                
                    except Exception as e:
                        print(f"Warning: Failed to extract image from slide {slide_num+1}: {str(e)}")
            
            # Add image content to summary
            if slide_images:
                summary_lines.append("\nIMAGES:")
                for i, img_info in enumerate(slide_images):
                    summary_lines.append(f"Image {i+1}:")
                    if img_info["caption"]:
                        summary_lines.append(f"  Description: {img_info['caption']}")
                    if img_info["ocr_text"]:
                        summary_lines.append(f"  Text content: {img_info['ocr_text']}")
            
            # If slide has no content, mark it
            if not slide_has_content:
                summary_lines.append("No content on this slide")
            
            # Add a separator between slides (double newline)
            summary_lines.append("\n")
    
    # Join all the summarized content
    combined_text = "\n".join(summary_lines)
    
    # Log the overall extraction result
    print(f"PowerPoint extraction complete: {len(combined_text)} characters extracted")
    if not combined_text.strip():
        print("WARNING: No content was extracted from the PowerPoint file!")
        
    return combined_text

# Initialize Flask app
app = Flask(__name__)
CORS(app)

# Configure OpenAI
openai.api_key = os.getenv("OPENAI_API_KEY")

# Configure Anthropic
anthropic_client = anthropic.Anthropic(api_key=os.getenv("ANTHROPIC_API_KEY", "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"))

# Configure Groq
GROQ_API_KEY = os.getenv("GROQ_API_KEY", "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")
GROQ_API_URL = "https://api.groq.com/openai/v1/chat/completions"

# Configure JWT
app.config['JWT_SECRET_KEY'] = os.getenv("JWT_SECRET", "default-dev-secret")
app.config['JWT_ACCESS_TOKEN_EXPIRES'] = datetime.timedelta(hours=24)

# Configure upload folders
UPLOAD_FOLDER = os.path.join(os.path.dirname(os.path.abspath(__file__)), "uploads")
DOCUMENT_FOLDER = os.path.join(UPLOAD_FOLDER, "documents")
IMAGE_FOLDER = os.path.join(UPLOAD_FOLDER, "images")

# Create folders if they don't exist
os.makedirs(DOCUMENT_FOLDER, exist_ok=True)
os.makedirs(IMAGE_FOLDER, exist_ok=True)

# Initialize ChromaDB
chroma_db_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "chroma_db")
os.makedirs(chroma_db_path, exist_ok=True)
chroma_client = chromadb.PersistentClient(path=chroma_db_path)
openai_ef = embedding_functions.OpenAIEmbeddingFunction(
    api_key=os.getenv("OPENAI_API_KEY"),
    model_name="text-embedding-ada-002"
)

# Initialize ImageProcessor for image RAG
app_base_dir = os.path.dirname(os.path.abspath(__file__))
image_processor = ImageProcessor(app_base_dir)

# Create directories for storing conversations
CONVERSATIONS_FOLDER = os.path.join(os.path.dirname(os.path.abspath(__file__)), "conversations")
os.makedirs(CONVERSATIONS_FOLDER, exist_ok=True)

# Sync datasets with ChromaDB collections
def sync_datasets_with_collections():
    print("Syncing datasets with ChromaDB collections...")
    datasets_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "datasets")
    os.makedirs(datasets_dir, exist_ok=True)
    
    # Get all dataset files
    dataset_files = [f for f in os.listdir(datasets_dir) if f.endswith('_datasets.json')]
    
    # In ChromaDB v0.6.0, list_collections() only returns collection names
    existing_collections = chroma_client.list_collections()
    
    for dataset_file in dataset_files:
        try:
            with open(os.path.join(datasets_dir, dataset_file), 'r') as f:
                datasets = json.load(f)
                
            for dataset in datasets:
                dataset_id = dataset["id"]
                
                # Check if collection exists in ChromaDB
                if dataset_id not in existing_collections:
                    print(f"Creating missing collection for dataset: {dataset_id}")
                    try:
                        chroma_client.create_collection(name=dataset_id, embedding_function=openai_ef)
                    except Exception as e:
                        print(f"Error creating collection for dataset {dataset_id}: {str(e)}")
        except Exception as e:
            print(f"Error processing dataset file {dataset_file}: {str(e)}")
    
    print("Dataset sync completed")

# Run the sync on app startup
sync_datasets_with_collections()

# Helper functions
def get_token_from_header():
    auth_header = request.headers.get('Authorization')
    if not auth_header or 'Bearer ' not in auth_header:
        return None
    return auth_header.split('Bearer ')[1]

def verify_token():
    token = get_token_from_header()
    if not token:
        return None
    try:
        decoded = jwt.decode(token, app.config['JWT_SECRET_KEY'], algorithms=["HS256"])
        return decoded
    except jwt.ExpiredSignatureError:
        return None
    except jwt.InvalidTokenError:
        return None

def require_auth(f):
    def decorated(*args, **kwargs):
        user_data = verify_token()
        if not user_data:
            return jsonify({"error": "Unauthorized"}), 401
        return f(user_data, *args, **kwargs)
    decorated.__name__ = f.__name__
    return decorated

def extract_text_from_file(file_path):
    extension = os.path.splitext(file_path)[1].lower()
    if extension == '.pdf':
        return extract_text_from_pdf(file_path)
    elif extension == '.docx':
        return extract_text_from_docx(file_path)
    elif extension == '.pptx':
        return extract_text_from_pptx(file_path)
    elif extension == '.txt':
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    else:
        return ""

def create_semantic_chunks(text, max_chunk_size=1000, overlap=200):
    """Create semantically coherent chunks from text, respecting document structure
    
    Args:
        text: The text to chunk
        max_chunk_size: Maximum size of a chunk in characters
        overlap: Number of characters to overlap between chunks
        
    Returns:
        list: List of text chunks
    """
    # Detect if content is from a PowerPoint presentation
    is_powerpoint = "## SLIDE " in text or "## PRESENTATION SUMMARY ##" in text
    
    # Use larger chunks for PowerPoint content
    if is_powerpoint:
        max_chunk_size = 3000  # Much larger to contain multiple slides
        overlap = 500  # More overlap to maintain context between chunks
    
    # If text is empty or too short, return as is
    if not text or len(text) <= max_chunk_size:
        return [text]
    
    # For PowerPoint content, use a special chunking strategy
    if is_powerpoint:
        return chunk_powerpoint_content(text, max_chunk_size, overlap)
    
    # Split the text into sections based on headers or meaningful separators
    # Look for common headers like "Chapter", "Section", or numbered headers
    section_patterns = [
        r'(?:\n|\r\n|\r)#{1,6}\s+(.+?)(?:\n|\r\n|\r)',  # Markdown headers
        r'(?:\n|\r\n|\r)([A-Z][^.\n]{0,50}:)(?:\n|\r\n|\r)',  # Capitalized section names with colon
        r'(?:\n|\r\n|\r)((?:Section|Chapter|Part)\s+\d+[.:]\s*[^\n]+)(?:\n|\r\n|\r)',  # Section/Chapter headings
        r'(?:\n|\r\n|\r)(\d+\.\d*\s+[^\n]{0,50})(?:\n|\r\n|\r)',  # Numbered sections (e.g., "1.2 Title")
        r'(?:\n|\r\n|\r)(Slide \d+:)',  # Slide markers
        r'(?:\n|\r\n|\r)(Title: .*?)(?:\n|\r\n|\r)',  # PowerPoint slide titles
        r'(?:\n|\r\n|\r)(Speaker Notes: .*?)(?:\n|\r\n|\r)',  # PowerPoint speaker notes
        r'(?:\n|\r\n|\r)(Image \d+\.\d+ text:)',  # Image text markers
        r'(?:\n|\r\n|\r)(Page \d+ image text:)',  # Page image text markers
        r'(?:\n|\r\n|\r)(Extracted from images:)',  # OCR section marker
        r'(?:\n|\r\n|\r)(Presentation Title: .*?)(?:\n|\r\n|\r)',  # Presentation title
        r'(?:\n|\r\n|\r)(Author: .*?)(?:\n|\r\n|\r)'  # Author info
    ]
    
    # Find all potential section boundaries
    sections = []
    last_end = 0
    
    # First, identify all possible section boundaries
    boundaries = []
    for pattern in section_patterns:
        for match in re.finditer(pattern, text, re.MULTILINE):
            boundaries.append((match.start(), match.group(0)))
    
    # Sort boundaries by position
    boundaries.sort(key=lambda x: x[0])
    
    # Create sections based on boundaries
    for i, (pos, header) in enumerate(boundaries):
        if pos > last_end:  # Only process if we haven't already included this position
            if i < len(boundaries) - 1:
                next_pos = boundaries[i+1][0]
                section_text = text[pos:next_pos]
                sections.append(section_text)
                last_end = next_pos
            else:
                # Last section
                section_text = text[pos:]
                sections.append(section_text)
                last_end = len(text)
    
    # If no sections were found or first section doesn't start at beginning
    if not sections or boundaries[0][0] > 0 if boundaries else True:
        if boundaries:
            first_pos = boundaries[0][0]
            sections.insert(0, text[:first_pos])
        else:
            # If no sections were identified, treat the whole text as one section
            sections = [text]
    
    # Further chunk each section if it's still too large
    chunks = []
    for section in sections:
        if len(section) <= max_chunk_size:
            chunks.append(section)
        else:
            # Try to split on paragraph boundaries
            paragraphs = re.split(r'(?:\n|\r\n|\r){2,}', section)
            current_chunk = ""
            
            for para in paragraphs:
                # If adding this paragraph exceeds max size and we already have content
                if len(current_chunk) + len(para) > max_chunk_size and current_chunk:
                    chunks.append(current_chunk)
                    # Start new chunk with overlap
                    overlap_point = max(0, len(current_chunk) - overlap)
                    current_chunk = current_chunk[overlap_point:] + para
                else:
                    current_chunk += para
                    
                if len(current_chunk) > max_chunk_size:
                    # If a single paragraph is too big, we have to split within the paragraph
                    sentences = re.split(r'(?<=[.!?])\s+', current_chunk)
                    current_chunk = ""
                    
                    for sentence in sentences:
                        if len(current_chunk) + len(sentence) > max_chunk_size and current_chunk:
                            chunks.append(current_chunk)
                            # Start new chunk with overlap
                            overlap_point = max(0, len(current_chunk) - overlap)
                            current_chunk = current_chunk[overlap_point:] + sentence
                        else:
                            current_chunk += sentence + " "
            
            # Don't forget to add the last chunk
            if current_chunk:
                chunks.append(current_chunk)
    
    # Post-processing to ensure chunks are not too small or too large
    final_chunks = []
    current_chunk = ""
    
    for chunk in chunks:
        # If chunk is very small, combine with previous unless it's a special section
        if len(chunk) < max_chunk_size / 4 and not any(re.search(pattern, chunk) for pattern in section_patterns):
            if current_chunk and len(current_chunk) + len(chunk) <= max_chunk_size:
                current_chunk += "\n\n" + chunk
            else:
                if current_chunk:
                    final_chunks.append(current_chunk)
                current_chunk = chunk
        else:
            if current_chunk:
                final_chunks.append(current_chunk)
            current_chunk = chunk
    
    # Add the last chunk
    if current_chunk:
        final_chunks.append(current_chunk)
    
    return final_chunks

def chunk_powerpoint_content(text, max_chunk_size=3000, overlap=500):
    """Special chunking function for PowerPoint content that keeps slides together
    
    Args:
        text: The PowerPoint text to chunk
        max_chunk_size: Maximum size of a chunk in characters
        overlap: Number of characters to overlap between chunks
        
    Returns:
        list: List of chunks with multiple slides grouped together
    """
    # If content is small enough, return as is
    if len(text) <= max_chunk_size:
        return [text]
    
    # Split the text by slide markers
    slides = re.split(r'(## SLIDE \d+ ##)', text)
    
    # Pair slide markers with their content
    slide_sections = []
    for i in range(0, len(slides) - 1, 2):
        if i + 1 < len(slides):
            slide_marker = slides[i]
            slide_content = slides[i + 1]
            slide_sections.append(slide_marker + slide_content)
        elif i < len(slides):
            # Handle odd number of splits
            slide_sections.append(slides[i])
    
    # If we didn't get proper slide sections, fallback to simple splitting
    if not slide_sections:
        slide_sections = [text]
    
    # Group slides into chunks
    chunks = []
    current_chunk = ""
    
    for slide in slide_sections:
        # If adding this slide would exceed the max size
        if len(current_chunk) + len(slide) > max_chunk_size and current_chunk:
            chunks.append(current_chunk)
            
            # Get the slides we just added to the chunk for overlap
            slide_markers_in_chunk = re.findall(r'## SLIDE (\d+) ##', current_chunk)
            
            # Start new chunk with presentation info and the last few slides for overlap
            # Aim to include about 2-3 slides for overlap if possible
            if slide_markers_in_chunk:
                # Extract presentation summary if it exists
                presentation_summary = ""
                if "## PRESENTATION SUMMARY ##" in text:
                    summary_match = re.search(r'(## PRESENTATION SUMMARY ##.*?)(?=## SLIDE)', text, re.DOTALL)
                    if summary_match:
                        presentation_summary = summary_match.group(1)
                
                # Find the last 2 slides to include for overlap
                overlap_slide_numbers = slide_markers_in_chunk[-2:] if len(slide_markers_in_chunk) >= 2 else slide_markers_in_chunk
                
                # Start new chunk with presentation summary and overlapping slides
                new_chunk_start = presentation_summary if presentation_summary else ""
                
                # Add overlapping slides
                for slide_num in overlap_slide_numbers:
                    slide_pattern = f"## SLIDE {slide_num} ##.*?(?=## SLIDE|$)"
                    slide_match = re.search(slide_pattern, text, re.DOTALL)
                    if slide_match:
                        new_chunk_start += slide_match.group(0)
                
                # Add the new slide
                current_chunk = new_chunk_start + slide
            else:
                # No slide markers found, use character-based overlap
                overlap_point = max(0, len(current_chunk) - overlap)
                current_chunk = current_chunk[overlap_point:] + slide
        else:
            current_chunk += slide
    
    # Add the last chunk
    if current_chunk:
        chunks.append(current_chunk)
    
    return chunks

# Add resize_image function before your routes
def resize_image(image_path, max_dimension=2048):
    """Resize image if it's too large while maintaining aspect ratio"""
    try:
        img = Image.open(image_path)
        
        # Check if resize is needed
        width, height = img.size
        if width <= max_dimension and height <= max_dimension:
            return image_path  # No resize needed
            
        # Calculate new dimensions
        if width > height:
            new_width = max_dimension
            new_height = int(height * (max_dimension / width))
        else:
            new_height = max_dimension
            new_width = int(width * (max_dimension / height))
            
        # Resize and save
        img = img.resize((new_width, new_height), Image.LANCZOS)
        
        # Create new filename for resized image
        filename, ext = os.path.splitext(image_path)
        resized_path = f"{filename}_resized{ext}"
        
        # Save resized image
        img.save(resized_path)
        return resized_path
    except Exception as e:
        print(f"Error resizing image: {str(e)}")
        return image_path  # Return original path if resizing fails

# Add the missing find_dataset_by_id function
def find_dataset_by_id(dataset_id):
    """Find a dataset by its ID across all users
    
    Args:
        dataset_id: The ID of the dataset to find
        
    Returns:
        dict: The dataset if found, None otherwise
    """
    datasets_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "datasets")
    if not os.path.exists(datasets_dir):
        return None
    
    # Check all user dataset files
    for filename in os.listdir(datasets_dir):
        if filename.endswith("_datasets.json"):
            try:
                with open(os.path.join(datasets_dir, filename), 'r') as f:
                    datasets = json.load(f)
                
                for dataset in datasets:
                    if dataset.get("id") == dataset_id:
                        return dataset
            except Exception as e:
                print(f"Error reading dataset file {filename}: {str(e)}")
    
    return None

# Authentication routes
@app.route('/api/auth/register', methods=['POST'])
def register():
    data = request.json
    username = data.get('username')
    password = data.get('password')
    
    if not username or not password:
        return jsonify({"error": "Username and password are required"}), 400
        
    # Check if user already exists - for demo purposes using a simple file
    users_file = os.path.join(os.path.dirname(os.path.abspath(__file__)), "users.json")
    users = {}
    if os.path.exists(users_file):
        with open(users_file, 'r') as f:
            users = json.load(f)
    
    if username in users:
        return jsonify({"error": "Username already exists"}), 400
        
    # Hash password and store user
    hashed_password = bcrypt.hashpw(password.encode('utf-8'), bcrypt.gensalt()).decode('utf-8')
    user_id = str(uuid.uuid4())
    users[username] = {
        "id": user_id,
        "username": username,
        "password": hashed_password
    }
    
    with open(users_file, 'w') as f:
        json.dump(users, f)
        
    # Generate token
    token = jwt.encode(
        {"id": user_id, "username": username, "exp": datetime.datetime.now(UTC) + app.config['JWT_ACCESS_TOKEN_EXPIRES']},
        app.config['JWT_SECRET_KEY']
    )
    
    return jsonify({"token": token, "user": {"id": user_id, "username": username}}), 201

@app.route('/api/auth/login', methods=['POST'])
def login():
    data = request.json
    username = data.get('username')
    password = data.get('password')
    
    if not username or not password:
        return jsonify({"error": "Username and password are required"}), 400
        
    # Load users - for demo purposes using a simple file
    users_file = os.path.join(os.path.dirname(os.path.abspath(__file__)), "users.json")
    if not os.path.exists(users_file):
        return jsonify({"error": "Invalid credentials"}), 401
        
    with open(users_file, 'r') as f:
        users = json.load(f)
        
    if username not in users:
        return jsonify({"error": "Invalid credentials"}), 401
        
    user = users[username]
    if not bcrypt.checkpw(password.encode('utf-8'), user["password"].encode('utf-8')):
        return jsonify({"error": "Invalid credentials"}), 401
        
    # Generate token
    token = jwt.encode(
        {"id": user["id"], "username": user["username"], "exp": datetime.datetime.now(UTC) + app.config['JWT_ACCESS_TOKEN_EXPIRES']},
        app.config['JWT_SECRET_KEY']
    )
    
    return jsonify({"token": token, "user": {"id": user["id"], "username": user["username"]}}), 200

# Dataset routes
@app.route('/api/datasets', methods=['GET'])
@require_auth
def get_datasets(user_data):
    datasets_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "datasets")
    os.makedirs(datasets_dir, exist_ok=True)
    user_datasets_file = os.path.join(datasets_dir, f"{user_data['id']}_datasets.json")
    
    if not os.path.exists(user_datasets_file):
        return jsonify([]), 200
        
    with open(user_datasets_file, 'r') as f:
        datasets = json.load(f)
    
    # Update dataset counts with accurate information
    for dataset in datasets:
        dataset_id = dataset.get("id")
        dataset_type = dataset.get("type", "text")
        
        # Initialize counts if they don't exist
        if "document_count" not in dataset:
            dataset["document_count"] = 0
        if "chunk_count" not in dataset:
            dataset["chunk_count"] = 0
        if "image_count" not in dataset:
            dataset["image_count"] = 0
        
        # Update text document count and chunk count if it's a text or mixed dataset
        if dataset_type in ["text", "mixed"]:
            try:
                collection = chroma_client.get_or_create_collection(
                    name=dataset_id,
                    embedding_function=openai_ef
                )
                # Get actual chunk count from ChromaDB
                chunk_count = collection.count()
                dataset["chunk_count"] = chunk_count
                
                # Calculate unique document count
                try:
                    # Get all document IDs
                    results = collection.get()
                    if results and results["metadatas"]:
                        # Extract unique document IDs
                        document_ids = set()
                        for metadata in results["metadatas"]:
                            if metadata and "document_id" in metadata:
                                document_ids.add(metadata["document_id"])
                        dataset["document_count"] = len(document_ids)
                except Exception as e:
                    print(f"Error calculating document count for dataset {dataset_id}: {str(e)}")
                    
            except Exception as e:
                print(f"Error getting document count for dataset {dataset_id}: {str(e)}")
        
        # Update image count if it's an image or mixed dataset
        if dataset_type in ["image", "mixed"]:
            try:
                # Check image processor for image counts
                image_count = 0
                if dataset_id in image_processor.image_metadata:
                    image_count = len(image_processor.image_metadata[dataset_id])
                dataset["image_count"] = image_count
                
                # Include image previews (first 3 images)
                if dataset_id in image_processor.image_metadata and image_count > 0:
                    # Get up to 3 images to display as previews
                    image_previews = []
                    for i, img_meta in enumerate(image_processor.image_metadata[dataset_id][:3]):
                        if 'path' in img_meta:
                            image_previews.append({
                                "id": img_meta.get("id", ""),
                                "url": f"/api/images/{os.path.basename(img_meta['path'])}",
                                "caption": img_meta.get("caption", "")
                            })
                    dataset["image_previews"] = image_previews
            except Exception as e:
                print(f"Error getting image count for dataset {dataset_id}: {str(e)}")
    
    return jsonify(datasets), 200

@app.route('/api/datasets', methods=['POST'])
@require_auth
def create_dataset(user_data):
    data = request.json
    name = data.get('name')
    description = data.get('description', '')
    dataset_type = data.get('type', 'mixed')  # Default to 'mixed' instead of 'text' to support both
    
    if not name:
        return jsonify({"error": "Dataset name is required"}), 400
        
    dataset_id = str(uuid.uuid4())
    
    # Save dataset metadata
    datasets_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "datasets")
    os.makedirs(datasets_dir, exist_ok=True)
    user_datasets_file = os.path.join(datasets_dir, f"{user_data['id']}_datasets.json")
    
    datasets = []
    if os.path.exists(user_datasets_file):
        with open(user_datasets_file, 'r') as f:
            datasets = json.load(f)
            
    dataset = {
        "id": dataset_id,
        "name": name,
        "description": description,
        "type": dataset_type,
        "user_id": user_data['id'],  # Store the user_id in the dataset
        "created_at": datetime.datetime.now(UTC).isoformat(),
        "document_count": 0,
        "image_count": 0  # Initialize image count to zero
    }
    datasets.append(dataset)
    
    with open(user_datasets_file, 'w') as f:
        json.dump(datasets, f)
        
    # Create collection in ChromaDB for text documents
    if dataset_type in ['text', 'mixed']:
        chroma_client.create_collection(name=dataset_id, embedding_function=openai_ef)
    
    # Initialize FAISS index for image datasets if needed
    if dataset_type in ['image', 'mixed']:
        # The image_processor will create the index when first image is added
        # But we can preemptively initialize it
        try:
            empty_index = faiss.IndexFlatIP(512)  # 512-dim for CLIP embeddings
            index_path = os.path.join(image_processor.indices_dir, f"{dataset_id}_index.faiss")
            faiss.write_index(empty_index, index_path)
            
            # Initialize empty metadata
            metadata_path = os.path.join(image_processor.indices_dir, f"{dataset_id}_metadata.json")
            with open(metadata_path, 'w') as f:
                json.dump([], f)
                
            # Initialize in-memory data structures
            image_processor.image_indices[dataset_id] = empty_index
            image_processor.image_metadata[dataset_id] = []
        except Exception as e:
            print(f"Warning: Error initializing image index for dataset {dataset_id}: {str(e)}")
            # Continue anyway since the index will be created on first image upload
    
    return jsonify(dataset), 201

@app.route('/api/datasets/<dataset_id>/documents', methods=['POST'])
@require_auth
def upload_document(user_data, dataset_id):
    if 'file' not in request.files:
        return jsonify({"error": "No file part"}), 400
        
    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "No selected file"}), 400
        
    # Get dataset information
    dataset = find_dataset_by_id(dataset_id)
    if not dataset:
        return jsonify({"error": "Dataset not found"}), 404
        
    # Check if dataset belongs to user or user is admin
    user_id = user_data.get('id')
    is_admin = user_data.get('isAdmin', False)
    
    # Print debug information
    print(f"Dataset user_id: {dataset.get('user_id')}, Current user_id: {user_id}, isAdmin: {is_admin}")
    
    # Handle legacy datasets that might not have user_id
    # 1. If dataset has no user_id, assume it belongs to the current user
    # 2. If dataset has user_id and it doesn't match current user (and user is not admin), deny access
    if 'user_id' in dataset and dataset['user_id'] != user_id and not is_admin:
        return jsonify({"error": "You don't have permission to access this dataset"}), 403
    
    # Check file extension
    filename = secure_filename(file.filename)
    file_extension = os.path.splitext(filename)[1].lower()
    
    # Get dataset type (text, image, or mixed)
    dataset_type = dataset.get("type", "text")
    
    # Check if file type is supported for this dataset
    if dataset_type == "text" and file_extension not in text_extensions:
        return jsonify({"error": f"Unsupported file type for text dataset. Supported types: {', '.join(text_extensions)}"}), 400
    elif dataset_type == "image" and file_extension not in image_extensions:
        return jsonify({"error": f"Unsupported file type for image dataset. Supported types: {', '.join(image_extensions)}"}), 400
    elif dataset_type == "mixed":
        if file_extension not in text_extensions and file_extension not in image_extensions:
            return jsonify({"error": f"Unsupported file type. Supported types: {', '.join(text_extensions + image_extensions)}"}), 400
    
    # If this is an image file and the dataset supports images
    if file_extension in image_extensions and dataset_type in ["image", "mixed"]:
        try:
            # Save the file with a unique name to prevent overwriting
            file_id = str(uuid.uuid4())
            original_filename = filename
            filename = f"{file_id}{file_extension}"
            file_path = os.path.join(UPLOAD_FOLDER, filename)
            file.save(file_path)
            
            # Add image to dataset
            img_processor = ImageProcessor(DATA_FOLDER)
            image_meta = img_processor.add_image_to_dataset(
                dataset_id, 
                file_path,
                {
                    "dataset_id": dataset_id,
                    "original_filename": original_filename,
                    "url": f"/uploads/{filename}",
                    "type": "image"
                }
            )
            
            # Update image count in dataset
            datasets_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "datasets")
            user_datasets_file = os.path.join(datasets_dir, f"{user_data['id']}_datasets.json")
            
            if os.path.exists(user_datasets_file):
                with open(user_datasets_file, 'r') as f:
                    datasets = json.load(f)
                
                for idx, ds in enumerate(datasets):
                    if ds["id"] == dataset_id:
                        if "image_count" in ds:
                            datasets[idx]["image_count"] += 1
                        else:
                            datasets[idx]["image_count"] = 1
                        break
                
                with open(user_datasets_file, 'w') as f:
                    json.dump(datasets, f)
            
            # Return success
            return jsonify({
                "id": image_meta.get("id", ""),
                "filename": original_filename,
                "type": "image",
                "caption": image_meta.get("caption", ""),
                "url": f"/uploads/{filename}"
            }), 201
            
        except Exception as e:
            if os.path.exists(file_path):
                os.remove(file_path)
            return jsonify({"error": f"Error processing image: {str(e)}"}), 500
    
    # If this is a document file and the dataset supports text
    elif file_extension in text_extensions and dataset_type in ["text", "mixed"]:
        # Save the file
        file_path = os.path.join(DOCUMENT_FOLDER, f"{uuid.uuid4()}_{filename}")
        file.save(file_path)
        
        # Extract text from file
        text = extract_text_from_file(file_path)
        if not text:
            os.remove(file_path)
            return jsonify({"error": "Could not extract text from file"}), 400
        
        # For PowerPoint files, clean up the text by removing unnecessary metadata
        if file_extension.lower() == '.pptx':
            # Remove any metadata lines at the beginning
            lines = text.split('\n')
            clean_lines = []
            skip_metadata = True
            
            for line in lines:
                # Skip metadata at the beginning (like Presentation Title, Author, etc.)
                if skip_metadata and (line.startswith("Presentation Title:") or 
                                     line.startswith("Author:") or 
                                     line.startswith("Subject:") or 
                                     line.startswith("Keywords:") or 
                                     line.startswith("Category:") or 
                                     line.startswith("Comments:")):
                    continue
                
                # Once we hit actual slide content, stop skipping
                if line.startswith("## SLIDE "):
                    skip_metadata = False
                
                # Skip speaker notes
                if "Speaker Notes:" in line:
                    continue
                
                clean_lines.append(line)
            
            # Rejoin the cleaned text
            text = '\n'.join(clean_lines)
        
        # Detect PowerPoint content for chunking strategy
        is_powerpoint = file_extension.lower() == '.pptx' or "## SLIDE " in text
        
        # Set chunking parameters based on file type
        if is_powerpoint:
            # For PowerPoint files, use large chunks to keep multiple slides together
            max_chunk_size = 3000
            overlap = 500
        else:
            # Default chunking parameters for other document types
            max_chunk_size = 1000  
            overlap = 200
        
        # Create chunks from text using the semantic chunking algorithm
        chunks = create_semantic_chunks(text, max_chunk_size=max_chunk_size, overlap=overlap)
        
        # Create document
        document_id = str(uuid.uuid4())
        
        # Add chunks to vector store
        try:
            chroma_collection = chroma_client.get_or_create_collection(
                name=dataset_id,
                embedding_function=openai_ef
            )
            
            # Create unique IDs for chunks
            chunk_ids = [f"{document_id}_{i}" for i in range(len(chunks))]
            
            # Add metadata to each chunk
            metadatas = [{
                "document_id": document_id,
                "dataset_id": dataset_id,
                "filename": filename,
                "source": filename,  # Add source field for compatibility
                "file_path": file_path,
                "chunk": i,
                "total_chunks": len(chunks),
                "file_type": file_extension.lower(),
                "is_powerpoint": is_powerpoint,  # Add flag to identify PowerPoint content
                "created_at": datetime.datetime.now(UTC).isoformat(),
            } for i in range(len(chunks))]
            
            # Add chunks to vector store
            chroma_collection.add(
                ids=chunk_ids,
                documents=chunks,
                metadatas=metadatas
            )
            
            # Update document count and chunk count in dataset
            datasets_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "datasets")
            user_datasets_file = os.path.join(datasets_dir, f"{user_data['id']}_datasets.json")
            
            if os.path.exists(user_datasets_file):
                with open(user_datasets_file, 'r') as f:
                    datasets = json.load(f)
                
                for idx, ds in enumerate(datasets):
                    if ds["id"] == dataset_id:
                        # Update document count
                        if "document_count" in ds:
                            datasets[idx]["document_count"] += 1
                        else:
                            datasets[idx]["document_count"] = 1
                            
                        # Update chunk count
                        if "chunk_count" in ds:
                            datasets[idx]["chunk_count"] += len(chunks)
                        else:
                            datasets[idx]["chunk_count"] = len(chunks)
                        break
                
                with open(user_datasets_file, 'w') as f:
                    json.dump(datasets, f)
            
            # Return success
            return jsonify({
                "id": document_id,
                "filename": filename,
                "type": "text",
                "chunks": len(chunks)
            }), 201
            
        except Exception as e:
            # Clean up the file if there was an error
            if os.path.exists(file_path):
                try:
                    os.remove(file_path)
                except:
                    pass
            return jsonify({"error": f"Error processing document: {str(e)}"}), 500
    else:
        return jsonify({"error": "Unsupported file type"}), 400

@app.route('/api/datasets/<dataset_id>/documents/<document_id>', methods=['DELETE'])
@require_auth
def remove_document(user_data, dataset_id, document_id):
    # Check if dataset exists and belongs to user
    datasets_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "datasets")
    user_datasets_file = os.path.join(datasets_dir, f"{user_data['id']}_datasets.json")
    
    if not os.path.exists(user_datasets_file):
        return jsonify({"error": "Dataset not found"}), 404
        
    with open(user_datasets_file, 'r') as f:
        datasets = json.load(f)
        
    dataset_exists = False
    dataset_index = -1
    for idx, dataset in enumerate(datasets):
        if dataset["id"] == dataset_id:
            dataset_exists = True
            dataset_index = idx
            break
            
    if not dataset_exists:
        return jsonify({"error": "Dataset not found"}), 404
    
    # Try to get the collection from ChromaDB
    try:
        collection = chroma_client.get_collection(name=dataset_id, embedding_function=openai_ef)
        
        # Get all chunks with the matching document_id in their metadata
        results = collection.get(
            where={"document_id": document_id}
        )
        
        if not results or len(results['ids']) == 0:
            return jsonify({"error": "Document not found in dataset"}), 404
        
        # Get number of chunks to remove for updating chunk count
        num_chunks_to_remove = len(results['ids'])
        
        # Delete the chunks from ChromaDB
        collection.delete(
            ids=results['ids']
        )
        
        # Update document count and chunk count in dataset
        if datasets[dataset_index]["document_count"] > 0:
            datasets[dataset_index]["document_count"] -= 1
            
        # Update chunk count if it exists
        if "chunk_count" in datasets[dataset_index]:
            datasets[dataset_index]["chunk_count"] -= num_chunks_to_remove
            if datasets[dataset_index]["chunk_count"] < 0:
                datasets[dataset_index]["chunk_count"] = 0
            
        with open(user_datasets_file, 'w') as f:
            json.dump(datasets, f)
            
        return jsonify({
            "message": "Document removed successfully",
            "chunks_removed": num_chunks_to_remove
        }), 200
    
    except Exception as e:
        print(f"Error removing document: {str(e)}")
        return jsonify({"error": f"Failed to remove document: {str(e)}"}), 500

# Admin routes
@app.route('/api/admin/rebuild_dataset/<dataset_id>', methods=['POST'])
@require_auth
def rebuild_dataset(user_data, dataset_id):
    # Check if dataset exists and belongs to user
    datasets_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "datasets")
    user_datasets_file = os.path.join(datasets_dir, f"{user_data['id']}_datasets.json")
    
    if not os.path.exists(user_datasets_file):
        return jsonify({"error": "Dataset not found"}), 404
        
    with open(user_datasets_file, 'r') as f:
        datasets = json.load(f)
        
    dataset_exists = False
    dataset_index = -1
    for idx, dataset in enumerate(datasets):
        if dataset["id"] == dataset_id:
            dataset_exists = True
            dataset_index = idx
            break
            
    if not dataset_exists:
        return jsonify({"error": "Dataset not found"}), 404
    
    # Reset document count
    datasets[dataset_index]["document_count"] = 0
    
    with open(user_datasets_file, 'w') as f:
        json.dump(datasets, f)
    
    # Try to delete the existing collection if it exists
    try:
        existing_collections = chroma_client.list_collections()
        if dataset_id in existing_collections:
            chroma_client.delete_collection(name=dataset_id)
            print(f"Deleted existing collection for dataset: {dataset_id}")
    except Exception as e:
        print(f"Error deleting collection: {str(e)}")
    
    # Create a new collection
    try:
        chroma_client.create_collection(name=dataset_id, embedding_function=openai_ef)
        return jsonify({"message": "Dataset collection has been rebuilt. Please re-upload your documents."}), 200
    except Exception as e:
        return jsonify({"error": f"Failed to rebuild dataset: {str(e)}"}), 500

@app.route('/api/datasets/<dataset_id>/status', methods=['GET'])
@require_auth
def dataset_status(user_data, dataset_id):
    # Check if dataset exists and belongs to user
    datasets_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "datasets")
    user_datasets_file = os.path.join(datasets_dir, f"{user_data['id']}_datasets.json")
    
    if not os.path.exists(user_datasets_file):
        return jsonify({"error": "Dataset not found"}), 404
        
    with open(user_datasets_file, 'r') as f:
        datasets = json.load(f)
        
    dataset = None
    for d in datasets:
        if d["id"] == dataset_id:
            dataset = d
            break
            
    if not dataset:
        return jsonify({"error": "Dataset not found"}), 404
    
    # Get dataset type
    dataset_type = dataset.get("type", "text")
    
    # Check ChromaDB collection status for text documents
    collection_exists = False
    doc_count = 0
    chunk_count = 0
    
    if dataset_type in ["text", "mixed"]:
        existing_collections = chroma_client.list_collections()
        collection_exists = dataset_id in existing_collections
        
        if collection_exists:
            try:
                collection = chroma_client.get_collection(name=dataset_id, embedding_function=openai_ef)
                chunk_count = collection.count()
                
                # Get document count by counting unique document_ids
                try:
                    results = collection.get()
                    if results and results["metadatas"]:
                        # Extract unique document IDs
                        document_ids = set()
                        for metadata in results["metadatas"]:
                            if metadata and "document_id" in metadata:
                                document_ids.add(metadata["document_id"])
                        doc_count = len(document_ids)
                except Exception as e:
                    print(f"Error calculating document count: {str(e)}")
                    # Fallback to existing document count in dataset
                    doc_count = dataset.get("document_count", 0)
                    
            except Exception as e:
                return jsonify({
                    "dataset": dataset,
                    "collection_exists": False,
                    "error": str(e)
                }), 200
    
    # Check image index status for images
    image_index_exists = False
    image_count = 0
    image_previews = []
    
    if dataset_type in ["image", "mixed"]:
        # Check if there's an image index for this dataset
        index_path = os.path.join(image_processor.indices_dir, f"{dataset_id}_index.faiss")
        image_index_exists = os.path.exists(index_path)
        
        # Get image count from image processor
        if dataset_id in image_processor.image_metadata:
            image_count = len(image_processor.image_metadata[dataset_id])
            
            # Include image previews (first 3 images)
            if image_count > 0:
                for i, img_meta in enumerate(image_processor.image_metadata[dataset_id][:3]):
                    if 'path' in img_meta:
                        image_previews.append({
                            "id": img_meta.get("id", ""),
                            "url": f"/api/images/{os.path.basename(img_meta['path'])}",
                            "caption": img_meta.get("caption", "")
                        })
    
    # Update dataset with accurate counts
    dataset["document_count"] = doc_count
    dataset["chunk_count"] = chunk_count
    dataset["image_count"] = image_count
    
    return jsonify({
        "dataset": dataset,
        "collection_exists": collection_exists,
        "document_count": doc_count,
        "chunk_count": chunk_count,
        "image_index_exists": image_index_exists,
        "image_count": image_count,
        "image_previews": image_previews
    }), 200

@app.route('/api/datasets/<dataset_id>/documents', methods=['GET'])
@require_auth
def get_dataset_documents(user_data, dataset_id):
    # Check if dataset exists and belongs to user
    datasets_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "datasets")
    user_datasets_file = os.path.join(datasets_dir, f"{user_data['id']}_datasets.json")
    
    if not os.path.exists(user_datasets_file):
        return jsonify({"error": "Dataset not found"}), 404
        
    with open(user_datasets_file, 'r') as f:
        datasets = json.load(f)
        
    dataset_exists = False
    for dataset in datasets:
        if dataset["id"] == dataset_id:
            dataset_exists = True
            break
            
    if not dataset_exists:
        return jsonify({"error": "Dataset not found"}), 404
    
    # Try to get the collection from ChromaDB
    try:
        collection = chroma_client.get_collection(name=dataset_id, embedding_function=openai_ef)
        
        # Query all documents in the collection
        results = collection.get()
        
        if not results or len(results['ids']) == 0:
            return jsonify({"documents": []}), 200
        
        # Process results to get unique documents with chunk counts
        documents = {}
        for i, metadata in enumerate(results['metadatas']):
            if metadata and 'document_id' in metadata and 'source' in metadata:
                doc_id = metadata['document_id']
                if doc_id not in documents:
                    documents[doc_id] = {
                        'id': doc_id,
                        'filename': metadata['source'],
                        'chunk_count': 1,
                        'file_type': metadata.get('file_type', ''),
                        'created_at': metadata.get('created_at', '')
                    }
                else:
                    documents[doc_id]['chunk_count'] += 1
        
        # Convert to list and sort by creation date (newest first)
        document_list = list(documents.values())
        document_list.sort(key=lambda x: x.get('created_at', ''), reverse=True)
        
        return jsonify({
            "documents": document_list,
            "total_documents": len(document_list),
            "total_chunks": len(results['ids'])
        }), 200
    
    except Exception as e:
        print(f"Error getting documents: {str(e)}")
        return jsonify({"error": f"Failed to get documents: {str(e)}"}), 500

@app.route('/api/datasets/<dataset_id>', methods=['DELETE'])
@require_auth
def delete_dataset(user_data, dataset_id):
    # Check if dataset exists and belongs to user
    datasets_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "datasets")
    user_datasets_file = os.path.join(datasets_dir, f"{user_data['id']}_datasets.json")
    
    if not os.path.exists(user_datasets_file):
        return jsonify({"error": "Dataset not found"}), 404
        
    with open(user_datasets_file, 'r') as f:
        datasets = json.load(f)
    
    # Find and remove the dataset
    dataset_index = -1
    dataset_type = "text"  # Default
    for idx, dataset in enumerate(datasets):
        if dataset["id"] == dataset_id:
            dataset_index = idx
            dataset_type = dataset.get("type", "text")  # Get dataset type for cleanup
            break
    
    if dataset_index == -1:
        return jsonify({"error": "Dataset not found"}), 404
    
    # Remove dataset from the list
    deleted_dataset = datasets.pop(dataset_index)
    
    # Save updated datasets list
    with open(user_datasets_file, 'w') as f:
        json.dump(datasets, f)
    
    # Delete ChromaDB collection if it exists and it's a text or mixed dataset
    if dataset_type in ["text", "mixed"]:
        try:
            existing_collections = chroma_client.list_collections()
            if dataset_id in existing_collections:
                chroma_client.delete_collection(name=dataset_id)
                print(f"Deleted ChromaDB collection for dataset: {dataset_id}")
        except Exception as e:
            print(f"Error deleting ChromaDB collection: {str(e)}")
            # Continue with deletion even if ChromaDB deletion fails
    
    # Delete image index and files if it exists and it's an image or mixed dataset
    if dataset_type in ["image", "mixed"]:
        try:
            # Get list of image files before deleting the dataset
            image_files_to_delete = []
            if dataset_id in image_processor.image_metadata:
                for img_meta in image_processor.image_metadata[dataset_id]:
                    if 'path' in img_meta and os.path.exists(img_meta['path']):
                        image_files_to_delete.append(img_meta['path'])
            
            # Use the image processor to delete the dataset
            image_processor.delete_dataset(dataset_id)
            print(f"Deleted image index for dataset: {dataset_id}")
            
            # Delete any actual image files associated with this dataset
            for img_path in image_files_to_delete:
                try:
                    os.remove(img_path)
                    print(f"Deleted image file: {img_path}")
                except Exception as e:
                    print(f"Error deleting image file {img_path}: {str(e)}")
        except Exception as e:
            print(f"Error deleting image index: {str(e)}")
            # Continue with deletion even if image index deletion fails
    
    # Check if any bots are using this dataset and notify in the response
    bots_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "bots")
    user_bots_file = os.path.join(bots_dir, f"{user_data['id']}_bots.json")
    affected_bots = []
    
    if os.path.exists(user_bots_file):
        try:
            with open(user_bots_file, 'r') as f:
                bots = json.load(f)
            
            # Check both old dataset_id field and new dataset_ids array
            for bot in bots:
                # Check old single dataset reference
                if bot.get("dataset_id") == dataset_id:
                    affected_bots.append(bot["name"])
                
                # Check new dataset_ids array 
                if dataset_id in bot.get("dataset_ids", []):
                    if bot["name"] not in affected_bots:  # Avoid duplicates
                        affected_bots.append(bot["name"])
                        
                        # Remove this dataset from the bot's dataset_ids
                        if "dataset_ids" in bot:
                            bot["dataset_ids"] = [d for d in bot["dataset_ids"] if d != dataset_id]
            
            # Update bots file if any were affected
            if affected_bots:
                with open(user_bots_file, 'w') as f:
                    json.dump(bots, f)
        except Exception as e:
            print(f"Error checking for affected bots: {str(e)}")
    
    return jsonify({
        "message": "Dataset deleted successfully", 
        "deleted_dataset": deleted_dataset["name"],
        "affected_bots": affected_bots
    }), 200

# Chat routes
@app.route('/api/bots/<bot_id>/chat', methods=['POST'])
@require_auth
def chat_with_bot(user_data, bot_id):
    data = request.json
    message = data.get('message', '')
    debug_mode = data.get('debug_mode', False)
    use_model_chorus = data.get('use_model_chorus', False)  # User's explicit choice to use model chorus
    chorus_id = data.get('chorus_id', '')  # A specific chorus ID to use
    conversation_id = data.get('conversation_id', '')  # The conversation this message belongs to
    
    # Check for image data in base64 format
    image_data = data.get('image_data', '')
    image_path = None
    has_image = False
    
    if not message and not image_data:
        return jsonify({"error": "Message or image is required"}), 400
    
    # Create a new conversation if no conversation_id provided
    if not conversation_id:
        conversation_id = str(uuid.uuid4())
    
    # If image is provided, save it to a temporary file
    if image_data:
        has_image = True
        try:
            # Extract the base64 content and file extension
            if ';base64,' in image_data:
                header, encoded = image_data.split(';base64,')
                file_ext = header.split('/')[-1]
            else:
                encoded = image_data
                file_ext = 'png'  # Default to PNG if not specified
            
            # Decode the base64 data
            decoded_image = base64.b64decode(encoded)
            
            # Save to a temporary file
            filename = f"chat_image_{str(uuid.uuid4())}.{file_ext}"
            image_path = os.path.join(IMAGE_FOLDER, filename)
            
            with open(image_path, 'wb') as f:
                f.write(decoded_image)
                
            # Resize image if needed
            image_path = resize_image(image_path)
            
            # Update message to include reference to the image
            if not message:
                message = "[Image uploaded]"
            else:
                message = f"{message} [Image uploaded]"
                
        except Exception as img_error:
            print(f"Error processing image: {str(img_error)}")
            return jsonify({"error": f"Failed to process image: {str(img_error)}"}), 400
    
    # Create a new message object
    user_message = {
        "id": str(uuid.uuid4()),
        "role": "user",
        "content": message,
        "timestamp": datetime.datetime.now(UTC).isoformat()
    }
    
    # If there's an image, add it to the message metadata
    if has_image:
        user_message["has_image"] = True
        user_message["image_path"] = image_path
    
    # Add message to conversation history
    conversation_file = os.path.join(CONVERSATIONS_FOLDER, f"{user_data['id']}_{bot_id}_{conversation_id}.json")
    conversation_exists = os.path.exists(conversation_file)
    
    if conversation_exists:
        # Load existing conversation
        with open(conversation_file, 'r') as f:
            conversation = json.load(f)
    else:
        # Create new conversation
        conversation = {
            "id": conversation_id,
            "bot_id": bot_id,
            "user_id": user_data['id'],
            "title": message[:40] + "..." if len(message) > 40 else message,  # Use first message as title
            "created_at": datetime.datetime.now(UTC).isoformat(),
            "updated_at": datetime.datetime.now(UTC).isoformat(),
            "messages": []
        }
    
    # Add user message to conversation
    conversation["messages"].append(user_message)
    conversation["updated_at"] = datetime.datetime.now(UTC).isoformat()
    
    # Save updated conversation
    with open(conversation_file, 'w') as f:
        json.dump(conversation, f)
        
    # Get bot info
    bots_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "bots")
    user_bots_file = os.path.join(bots_dir, f"{user_data['id']}_bots.json")
    
    if not os.path.exists(user_bots_file):
        return jsonify({"error": "Bot not found"}), 404
        
    with open(user_bots_file, 'r') as f:
        bots = json.load(f)
        
    bot = None
    for b in bots:
        if b["id"] == bot_id:
            bot = b
            break
            
    if not bot:
        return jsonify({"error": "Bot not found"}), 404
        
    # Get dataset IDs from the bot
    dataset_ids = bot.get("dataset_ids", [])
    
    # Process image with OpenAI Vision API if an image is present
    if has_image:
        try:
            # Read the image file and encode it as base64
            with open(image_path, 'rb') as img_file:
                encoded_image = base64.b64encode(img_file.read()).decode('utf-8')
            
            # Prepare the message with the image for OpenAI Vision API
            messages = [
                {
                    "role": "system", 
                    "content": bot.get("system_instruction", "You are a helpful assistant that can analyze images.")
                },
                {
                    "role": "user", 
                    "content": [
                        {"type": "text", "text": message if message != "[Image uploaded]" else "What's in this image?"},
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/{file_ext};base64,{encoded_image}"
                            }
                        }
                    ]
                }
            ]
            
            # Call OpenAI API with vision capabilities
            response = openai.chat.completions.create(
                model="gpt-4o",
                messages=messages,
                max_tokens=1024
            )
            
            response_text = response.choices[0].message.content
            
            # Save the response in the conversation
            bot_response = {
                "id": str(uuid.uuid4()),
                "role": "assistant",
                "content": response_text,
                "timestamp": datetime.datetime.now(UTC).isoformat(),
                "from_image_analysis": True
            }
            conversation["messages"].append(bot_response)
            with open(conversation_file, 'w') as f:
                json.dump(conversation, f)
                
            return jsonify({
                "response": response_text,
                "conversation_id": conversation_id,
                "image_processed": True
            }), 200
            
        except Exception as vision_error:
            print(f"Error processing image with Vision API: {str(vision_error)}")
            # If the vision API fails, continue with the regular RAG process
    
    if not dataset_ids:
        return jsonify({
            "response": "I don't have any datasets to work with. Please add a dataset to help me answer your questions.",
            "conversation_id": conversation_id
        }), 200
    
    # Check if the bot has a chorus configuration associated with it
    # If so, use model chorus by default unless explicitly turned off
    bot_has_chorus = bot.get("chorus_id", "")
    use_model_chorus = use_model_chorus or bool(bot_has_chorus)
    
    # If a specific chorus_id is provided, use it; otherwise use the bot's chorus_id if available
    specific_chorus_id = chorus_id or bot.get("chorus_id", "")
    
    # Retrieve relevant documents from all datasets
    all_contexts = []
    image_results = []
    
    try:
        for dataset_id in dataset_ids:
            try:
                # Text-based document retrieval
                collection = chroma_client.get_collection(name=dataset_id, embedding_function=openai_ef)
                
                # Check if collection has any documents
                collection_count = collection.count()
                if collection_count > 0:
                    # Determine how many results to request based on collection size
                    n_results = min(5, collection_count)  # Default to 5 results
                    
                    # First get a sample to check if we're dealing with PowerPoint content
                    sample_results = collection.query(
                        query_texts=[message],
                        n_results=1
                    )
                    
                    # Check if the sample contains PowerPoint content
                    is_powerpoint = False
                    if sample_results and sample_results["documents"] and sample_results["documents"][0]:
                        sample_text = sample_results["documents"][0][0]
                        if "## SLIDE " in sample_text or "## PRESENTATION SUMMARY ##" in sample_text:
                            is_powerpoint = True
                            # For PowerPoint content, get more chunks to ensure we have enough context
                            n_results = min(10, collection_count)  # Double the number of chunks for PowerPoint
                    
                    # Get the actual results
                    results = collection.query(
                        query_texts=[message],
                        n_results=n_results
                    )
                    
                    all_contexts.extend(results["documents"][0])
                
                # Image-based retrieval if available
                try:
                    # First, check if dataset contains images
                    has_images = False
                    dataset_type = "text"  # Default
                    
                    # Check dataset type and image count from metadata
                    datasets_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "datasets")
                    for user_id in [user_data['id'], "shared"]:  # Check both user and shared datasets
                        user_datasets_file = os.path.join(datasets_dir, f"{user_id}_datasets.json")
                        if os.path.exists(user_datasets_file):
                            with open(user_datasets_file, 'r') as f:
                                datasets = json.load(f)
                                
                                for ds in datasets:
                                    if ds["id"] == dataset_id:
                                        dataset_type = ds.get("type", "text")
                                        if dataset_type in ["image", "mixed"] and ds.get("image_count", 0) > 0:
                                            has_images = True
                                        break
                    
                    # If dataset has images, retrieve them regardless of query type
                    if has_images:
                        # Search for images with the user's query
                        img_results = image_processor.search_images(dataset_id, message, top_k=3)
                        
                        # If no results or weak results, also try a more generic search
                        if not img_results or (img_results and img_results[0].get("score", 0) < 0.2):
                            # Try searching with a more generic query 
                            generic_queries = [
                                "relevant image for this topic",
                                "visual representation",
                                "image related to this subject"
                            ]
                            
                            for generic_query in generic_queries:
                                generic_results = image_processor.search_images(dataset_id, generic_query, top_k=2)
                                if generic_results:
                                    for gen_img in generic_results:
                                        # Add to results if not already included
                                        if not any(img.get("id") == gen_img.get("id") for img in img_results):
                                            img_results.append(gen_img)
                                    break  # Stop after finding some results
                        
                        # Add relevant images to results
                        if img_results:
                            for img in img_results:
                                # Add image information to results
                                image_info = {
                                    "id": img["id"],
                                    "caption": img["caption"],
                                    "url": f"/api/images/{os.path.basename(img['path'])}",
                                    "score": img["score"],
                                    "dataset_id": dataset_id
                                }
                                image_results.append(image_info)
                except Exception as img_error:
                    print(f"Error with image retrieval for dataset '{dataset_id}': {str(img_error)}")
                    # Continue even if image retrieval fails
                
            except Exception as coll_error:
                print(f"Error with collection '{dataset_id}': {str(coll_error)}")
                # Continue with other datasets even if one fails
                continue
                
        if not all_contexts and not image_results:
            # Save the response in the conversation
            bot_response = {
                "id": str(uuid.uuid4()),
                "role": "assistant",
                "content": "I don't have any documents or images in my knowledge base yet. Please upload some content to help me answer your questions.",
                "timestamp": datetime.datetime.now(UTC).isoformat()
            }
            conversation["messages"].append(bot_response)
            with open(conversation_file, 'w') as f:
                json.dump(conversation, f)
                
            return jsonify({
                "response": "I don't have any documents or images in my knowledge base yet. Please upload some content to help me answer your questions.",
                "debug": {"error": "No contexts found in any collections"} if debug_mode else None,
                "conversation_id": conversation_id
            }), 200
            
        # Sort contexts by relevance (they should already be sorted from query results)
        # But we need to truncate to avoid token limits
        
        # Check if we have PowerPoint content in the retrieved contexts
        has_powerpoint = any("## SLIDE " in ctx or "## PRESENTATION SUMMARY ##" in ctx for ctx in all_contexts)
        
        # Adjust the maximum number of contexts based on content type
        if has_powerpoint:
            max_contexts = 12  # Allow more contexts for PowerPoint content
        else:
            max_contexts = 8  # Default limit for regular content
            
        contexts = all_contexts[:max_contexts]
        context_text = "\n\n".join([f"[{i+1}] {ctx}" for i, ctx in enumerate(contexts)])
        
        # Prepare image information for context
        image_context = ""
        if image_results:
            # Sort images by score
            image_results.sort(key=lambda x: x["score"], reverse=True)
            
            # Take top images up to a limit
            top_images = image_results[:3]  # Limit to 3 images to avoid overloading context
            
            # Format image information
            image_context = "\n\nRelevant Images:\n"
            for i, img in enumerate(top_images):
                image_context += f"[Image {i+1}] Caption: {img['caption']}\n"
                image_context += f"            This image is available for viewing and download.\n"
            
            # Add the URLs to image metadata for later
            top_image_urls = [img["url"] for img in top_images]
            
        # Combine text and image contexts
        if context_text and image_context:
            full_context = context_text + image_context
        elif context_text:
            full_context = context_text
        elif image_context:
            full_context = image_context
        else:
            full_context = "No relevant information found."
        
        # Get source documents information with indices
        source_documents = []
        seen_documents = set()
        source_images = []
        
        # Add text document sources if available
        if all_contexts:
            for i, (result, metadata) in enumerate(zip(results["documents"][0][:max_contexts], results["metadatas"][0][:max_contexts])):
                if metadata and 'document_id' in metadata and 'source' in metadata:
                    doc_key = f"{metadata['document_id']}_{metadata['source']}"
                    if doc_key not in seen_documents:
                        source_documents.append({
                            'id': metadata['document_id'],
                            'filename': metadata['source'],
                            'dataset_id': dataset_id,
                            'context_index': i + 1  # Store the index used in context_text
                        })
                        seen_documents.add(doc_key)
        
        # Add image sources if available
        if image_results:
            for i, img in enumerate(top_images):
                source_images.append({
                    'id': img['id'],
                    'caption': img['caption'],
                    'url': img['url'],
                    'dataset_id': img['dataset_id'],
                    'type': 'image',
                    'context_index': f"Image {i+1}"
                })
        
        # Get previous conversation messages to add as context
        conversation_history = ""
        if len(conversation["messages"]) > 1:  # If there's more than just the current message
            # Get last 10 messages maximum
            recent_messages = conversation["messages"][-10:] if len(conversation["messages"]) > 10 else conversation["messages"]
            # Format them for context
            conversation_history = "\n".join([
                f"{msg['role'].capitalize()}: {msg['content']}" 
                for msg in recent_messages[:-1]  # Exclude current message
            ])
        
        # Prepare the system instruction and context
        system_instruction = bot.get("system_instruction", "You are a helpful assistant that answers questions based on the provided context.")
        
        # Add specific instructions for handling PowerPoint content
        system_instruction += "\n\nWhen answering questions about PowerPoint presentations, pay particular attention to slide numbers, titles, and slide content structure. Information in slide titles and speaker notes is particularly important. If information appears to be missing or incomplete in the retrieved context, explain that, but still try to provide the best possible answer."
        
        # Add conversation history to system instruction if available
        system_instruction_with_history = system_instruction
        if conversation_history:
            system_instruction_with_history += "\n\nThis is the conversation history so far:\n" + conversation_history
        
        # Check if using model chorus
        if use_model_chorus:
            # Check if a chorus configuration exists
            chorus_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "chorus")
            os.makedirs(chorus_dir, exist_ok=True)
            
            # Get the user's chorus definitions file
            user_choruses_file = os.path.join(chorus_dir, f"{user_data['id']}_choruses.json")
            
            # Try to load the specific chorus from the user's choruses
            chorus_config = None
            
            if os.path.exists(user_choruses_file) and specific_chorus_id:
                try:
                    with open(user_choruses_file, 'r') as f:
                        choruses = json.load(f)
                        
                    # Find the requested chorus
                    chorus_config = next((c for c in choruses if c["id"] == specific_chorus_id), None)
                    
                    if debug_mode and chorus_config:
                        print(f"Using chorus configuration: {chorus_config.get('name', 'Unnamed')}")
                except Exception as e:
                    print(f"Error loading chorus configuration: {str(e)}")
            
            # If no chorus config found, return an error
            if not chorus_config:
                # Save the response in the conversation
                bot_response = {
                    "id": str(uuid.uuid4()),
                    "role": "assistant",
                    "content": "I couldn't find the model chorus configuration. Please check that the chorus exists and is properly configured.",
                    "timestamp": datetime.datetime.now(UTC).isoformat()
                }
                conversation["messages"].append(bot_response)
                with open(conversation_file, 'w') as f:
                    json.dump(conversation, f)
                    
                return jsonify({
                    "response": "I couldn't find the model chorus configuration. Please check that the chorus exists and is properly configured.",
                    "debug": {"error": "Chorus configuration not found"} if debug_mode else None,
                    "conversation_id": conversation_id
                }), 200
            
            # Get response and evaluator models
            response_models = chorus_config.get('response_models', [])
            evaluator_models = chorus_config.get('evaluator_models', [])
            
            # For debugging
            if debug_mode:
                print(f"Using chorus configuration: {chorus_config.get('name', 'Unnamed')}")
                print(f"Response models: {len(response_models)}")
                print(f"Evaluator models: {len(evaluator_models)}")
            
            # Validate the configuration
            if not response_models or not evaluator_models:
                # Save the response in the conversation
                bot_response = {
                    "id": str(uuid.uuid4()),
                    "role": "assistant",
                    "content": "The model chorus configuration is incomplete. Please configure both response and evaluator models.",
                    "timestamp": datetime.datetime.now(UTC).isoformat()
                }
                conversation["messages"].append(bot_response)
                with open(conversation_file, 'w') as f:
                    json.dump(conversation, f)
                    
                return jsonify({
                    "response": "The model chorus configuration is incomplete. Please configure both response and evaluator models.",
                    "debug": {"error": "Incomplete chorus configuration"} if debug_mode else None,
                    "conversation_id": conversation_id
                }), 200
            
            logs = []
            logs.append(f"Using model chorus: {chorus_config.get('name', 'Unnamed')}")
            
            # Generate responses from all models
            all_responses = []
            anonymized_responses = []
            response_metadata = []
            
            # Process response models
            for model in response_models:
                provider = model.get('provider')
                model_name = model.get('model')
                temperature = float(model.get('temperature', 0.7))
                weight = int(model.get('weight', 1))
                
                for i in range(weight):
                    try:
                        if provider == 'OpenAI':
                            response = openai.chat.completions.create(
                                model=model_name,
                                messages=[
                                    {"role": "system", "content": system_instruction_with_history},
                                    {"role": "user", "content": f"Context:\n{context_text}\n\nUser question: {message}"}
                                ],
                                temperature=temperature
                            )
                            response_text = response.choices[0].message.content
                            anonymized_responses.append(response_text)
                            response_metadata.append({
                                "provider": provider,
                                "model": model_name,
                                "temperature": temperature
                            })
                            all_responses.append({
                                "provider": provider,
                                "model": model_name,
                                "response": response_text,
                                "temperature": temperature
                            })
                        elif provider == 'Anthropic':
                            response = anthropic_client.messages.create(
                                model=model_name,
                                system=system_instruction_with_history,
                                messages=[{"role": "user", "content": f"Context:\n{context_text}\n\nUser question: {message}"}],
                                temperature=temperature,
                                max_tokens=1024
                            )
                            response_text = response.content[0].text
                            anonymized_responses.append(response_text)
                            response_metadata.append({
                                "provider": provider,
                                "model": model_name,
                                "temperature": temperature
                            })
                            all_responses.append({
                                "provider": provider,
                                "model": model_name,
                                "response": response_text,
                                "temperature": temperature
                            })
                        elif provider == 'Groq':
                            headers = {
                                "Authorization": f"Bearer {GROQ_API_KEY}",
                                "Content-Type": "application/json"
                            }
                            payload = {
                                "model": model_name,
                                "messages": [
                                    {"role": "system", "content": system_instruction_with_history},
                                    {"role": "user", "content": f"Context:\n{context_text}\n\nUser question: {message}"}
                                ],
                                "temperature": temperature,
                                "max_tokens": 1024
                            }
                            response = requests.post(GROQ_API_URL, json=payload, headers=headers)
                            response_json = response.json()
                            response_text = response_json["choices"][0]["message"]["content"]
                            anonymized_responses.append(response_text)
                            response_metadata.append({
                                "provider": provider,
                                "model": model_name,
                                "temperature": temperature
                            })
                            all_responses.append({
                                "provider": provider,
                                "model": model_name,
                                "response": response_text,
                                "temperature": temperature
                            })
                        elif provider == 'Mistral':
                            # Note: This would require a Mistral API implementation
                            # For now we'll use OpenAI as a fallback and log it
                            logs.append(f"Mistral API not implemented, using OpenAI fallback for {model_name}")
                            response = openai.chat.completions.create(
                                model="gpt-3.5-turbo",
                                messages=[
                                    {"role": "system", "content": system_instruction_with_history},
                                    {"role": "user", "content": f"Context:\n{context_text}\n\nUser question: {message}"}
                                ],
                                temperature=temperature
                            )
                            response_text = response.choices[0].message.content
                            anonymized_responses.append(response_text)
                            response_metadata.append({
                                "provider": "OpenAI (Mistral fallback)",
                                "model": "gpt-3.5-turbo",
                                "temperature": temperature
                            })
                            all_responses.append({
                                "provider": "OpenAI (Mistral fallback)",
                                "model": "gpt-3.5-turbo",
                                "response": response_text,
                                "temperature": temperature
                            })
                    except Exception as e:
                        logs.append(f"Error with {provider} {model_name}: {str(e)}")
            
            # If no responses, use fallback
            if not all_responses:
                # Save the response in the conversation
                bot_response = {
                    "id": str(uuid.uuid4()),
                    "role": "assistant",
                    "content": "I encountered an issue with the model chorus. No models were able to generate a response.",
                    "timestamp": datetime.datetime.now(UTC).isoformat()
                }
                conversation["messages"].append(bot_response)
                with open(conversation_file, 'w') as f:
                    json.dump(conversation, f)
                    
                return jsonify({
                    "response": "I encountered an issue with the model chorus. No models were able to generate a response.",
                    "debug": {"error": "No models returned responses", "logs": logs} if debug_mode else None,
                    "conversation_id": conversation_id
                }), 200
            
            # If only one response, return it directly
            if len(all_responses) == 1:
                response_text = all_responses[0]["response"]
                
                # Save the response in the conversation
                bot_response = {
                    "id": str(uuid.uuid4()),
                    "role": "assistant",
                    "content": response_text,
                    "timestamp": datetime.datetime.now(UTC).isoformat()
                }
                conversation["messages"].append(bot_response)
                with open(conversation_file, 'w') as f:
                    json.dump(conversation, f)
                    
                return jsonify({
                    "response": response_text,
                    "debug": {
                        "all_responses": all_responses,
                        "anonymized_responses": anonymized_responses,
                        "response_metadata": response_metadata,
                        "logs": logs,
                        "contexts": contexts
                    } if debug_mode else None,
                    "conversation_id": conversation_id
                }), 200
            
            # Get voting from evaluator models
            votes = [0] * len(all_responses)
            
            for model in evaluator_models:
                provider = model.get('provider')
                model_name = model.get('model')
                temperature = float(model.get('temperature', 0.2))
                weight = int(model.get('weight', 1))
                
                voting_prompt = f"""You are an expert AI response evaluator. You need to rank the following responses to the question: "{message}"

Here are the {len(anonymized_responses)} candidate responses:

{chr(10).join([f"Response {j+1}:\n{resp}" for j, resp in enumerate(anonymized_responses)])}

Which response provides the most accurate, helpful, and relevant answer? Return ONLY the number (1-{len(anonymized_responses)}) of the best response.
Do not reveal any bias or preference based on writing style or approach - evaluate solely on answer quality, accuracy and helpfulness.
"""
                # Apply weight by counting vote multiple times
                for i in range(weight):
                    try:
                        vote_text = ""
                        if provider == 'OpenAI':
                            voting_response = openai.chat.completions.create(
                                model=model_name,
                                messages=[{"role": "user", "content": voting_prompt}],
                                temperature=temperature
                            )
                            vote_text = voting_response.choices[0].message.content
                        elif provider == 'Anthropic':
                            voting_response = anthropic_client.messages.create(
                                model=model_name,
                                messages=[{"role": "user", "content": voting_prompt}],
                                temperature=temperature,
                                max_tokens=10
                            )
                            vote_text = voting_response.content[0].text
                        elif provider == 'Groq':
                            headers = {
                                "Authorization": f"Bearer {GROQ_API_KEY}",
                                "Content-Type": "application/json"
                            }
                            payload = {
                                "model": model_name,
                                "messages": [{"role": "user", "content": voting_prompt}],
                                "temperature": temperature,
                                "max_tokens": 10
                            }
                            voting_response = requests.post(GROQ_API_URL, json=payload, headers=headers)
                            voting_json = voting_response.json()
                            vote_text = voting_json["choices"][0]["message"]["content"]
                        elif provider == 'Mistral':
                            # Fallback to OpenAI for Mistral
                            logs.append(f"Mistral API not implemented for evaluation, using OpenAI fallback")
                            voting_response = openai.chat.completions.create(
                                model="gpt-3.5-turbo",
                                messages=[{"role": "user", "content": voting_prompt}],
                                temperature=temperature
                            )
                            vote_text = voting_response.choices[0].message.content
                        
                        # Extract vote number
                        match = re.search(r'\b([1-9][0-9]*)\b', vote_text)
                        if match:
                            vote_number = int(match.group(1))
                            if 1 <= vote_number <= len(anonymized_responses):
                                votes[vote_number-1] += 1
                                logs.append(f"Evaluator {provider} {model_name} voted for response {vote_number}")
                            else:
                                logs.append(f"Invalid vote number: {vote_number}")
                        else:
                            logs.append(f"Could not parse vote from: {vote_text}")
                    except Exception as e:
                        logs.append(f"Error with evaluator {provider} {model_name}: {str(e)}")
            
            # Find winning response
            max_votes = max(votes) if votes else 0
            winning_indices = [i for i, v in enumerate(votes) if v == max_votes]
            winning_index = winning_indices[0] if winning_indices else 0
            
            winning_response = all_responses[winning_index]["response"]
            
            # Save the response in the conversation
            bot_response = {
                "id": str(uuid.uuid4()),
                "role": "assistant",
                "content": winning_response,
                "timestamp": datetime.datetime.now(UTC).isoformat()
            }
            conversation["messages"].append(bot_response)
            with open(conversation_file, 'w') as f:
                json.dump(conversation, f)
                
            return jsonify({
                "response": winning_response,
                "source_documents": source_documents,
                "debug": {
                    "all_responses": all_responses,
                    "anonymized_responses": anonymized_responses,
                    "response_metadata": response_metadata,
                    "votes": votes,
                    "logs": logs,
                    "contexts": contexts
                } if debug_mode else None,
                "conversation_id": conversation_id
            }), 200
                
    
        # Standard mode - just get a response from OpenAI
        response = openai.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": system_instruction_with_history + "\n\nWhen you reference information from the provided context, please cite the source using the number in square brackets, e.g. [1], [2], etc. For images, use [Image 1], [Image 2], etc. When referencing images, mention that users can click or download the image for a better view."},
                {"role": "user", "content": f"Context:\n{full_context}\n\nUser question: {message}\n\nPlease include citations [1], [2], etc. when you reference information from the context."}
            ]
        )
        
        response_text = response.choices[0].message.content
        
        # Extract which context indices were actually used in the response
        used_indices = set()
        for i in range(1, max_contexts + 1):
            if f"[{i}]" in response_text:
                used_indices.add(i)
        
        # Extract which image indices were referenced
        used_image_indices = set()
        for i in range(1, len(top_images) + 1) if image_results else []:
            if f"[Image {i}]" in response_text:
                used_image_indices.add(i)
        
        # Filter source_documents to only include those that were actually cited
        used_source_documents = [doc for doc in source_documents if doc['context_index'] in used_indices]
        
        # Filter source_images to only include those that were actually cited
        used_source_images = []
        if image_results:
            for img in source_images:
                img_idx = int(img['context_index'].split()[1])
                if img_idx in used_image_indices:
                    used_source_images.append(img)

        # Create detailed context information for frontend display
        context_details = []
        for i, ctx in enumerate(contexts):
            if i+1 in used_indices:  # Only include contexts that were cited
                # Find the corresponding source document
                source_doc = next((doc for doc in source_documents if doc['context_index'] == i+1), None)
                if source_doc:
                    context_details.append({
                        "index": i+1,
                        "snippet": ctx[:200] + "..." if len(ctx) > 200 else ctx,  # Truncate for display
                        "full_text": ctx,  # Include full snippet text
                        "document_id": source_doc.get('id', ''),
                        "filename": source_doc.get('filename', ''),
                        "dataset_id": source_doc.get('dataset_id', '')
                    })

        # Add image details to context details for referenced images
        image_details = []
        if used_source_images:
            for img in used_source_images:
                # Create specific download URL by adding download=true parameter
                base_url = img['url']
                download_url = f"{base_url}?download=true"
                
                image_details.append({
                    "index": img['context_index'],
                    "caption": img['caption'],
                    "url": img['url'],
                    "download_url": download_url,
                    "id": img['id'],
                    "dataset_id": img['dataset_id']
                })

        # Save the response in the conversation
        bot_response = {
            "id": str(uuid.uuid4()),
            "role": "assistant",
            "content": response_text,
            "timestamp": datetime.datetime.now(UTC).isoformat(),
            "referenced_images": [img["url"] for img in used_source_images] if used_source_images else [],
            "context_details": context_details,  # Add context details to the saved response
            "image_details": image_details  # Add image details to the saved response
        }
        conversation["messages"].append(bot_response)
        with open(conversation_file, 'w') as f:
            json.dump(conversation, f)
            
        return jsonify({
            "response": response_text,
            "source_documents": used_source_documents,
            "source_images": used_source_images,
            "context_details": context_details,  # Include context details in response
            "image_details": image_details,  # Include image details in response
            "debug": {
                "contexts": contexts,
                "image_results": image_results if image_results else []
            } if debug_mode else None,
            "conversation_id": conversation_id
        }), 200
        
    except Exception as e:
        print(f"Error in chat_with_bot: {str(e)}")
        # Save the error response in the conversation
        bot_response = {
            "id": str(uuid.uuid4()),
            "role": "assistant",
            "content": "I apologize, but I encountered an error while processing your request. Please try again or contact support if the issue persists.",
            "timestamp": datetime.datetime.now(UTC).isoformat()
        }
        conversation["messages"].append(bot_response)
        with open(conversation_file, 'w') as f:
            json.dump(conversation, f)
            
        return jsonify({
            "response": "I apologize, but I encountered an error while processing your request. Please try again or contact support if the issue persists.",
            "debug": {"error": str(e)} if debug_mode else None,
            "conversation_id": conversation_id
        }), 200

# New endpoints for conversation management
@app.route('/api/bots/<bot_id>/conversations', methods=['GET'])
@require_auth
def get_conversations(user_data, bot_id):
    """Get all conversations for a specific bot"""
    try:
        # Read all conversation files for this user and bot
        conversation_files = [f for f in os.listdir(CONVERSATIONS_FOLDER) 
                             if f.startswith(f"{user_data['id']}_{bot_id}_") and f.endswith('.json')]
        
        conversations = []
        for file_name in conversation_files:
            file_path = os.path.join(CONVERSATIONS_FOLDER, file_name)
            try:
                with open(file_path, 'r') as f:
                    conversation = json.load(f)
                
                # Add a preview of the conversation
                if conversation.get("messages"):
                    # Get the first user message as title if not already set
                    if not conversation.get("title"):
                        first_user_message = next((msg for msg in conversation["messages"] if msg["role"] == "user"), None)
                        if first_user_message:
                            title = first_user_message["content"]
                            conversation["title"] = title[:40] + "..." if len(title) > 40 else title
                    
                    # Count messages
                    message_count = len(conversation["messages"])
                    conversation["message_count"] = message_count
                
                # Add a summary of the conversation
                conversations.append({
                    "id": conversation["id"],
                    "title": conversation.get("title", "Untitled Conversation"),
                    "created_at": conversation.get("created_at"),
                    "updated_at": conversation.get("updated_at"),
                    "message_count": conversation.get("message_count", 0),
                    "preview": conversation["messages"][-1]["content"][:100] + "..." if conversation.get("messages") else ""
                })
            except Exception as e:
                print(f"Error loading conversation {file_name}: {str(e)}")
        
        # Sort conversations by updated_at (newest first)
        conversations.sort(key=lambda x: x.get("updated_at", ""), reverse=True)
        
        return jsonify(conversations), 200
    except Exception as e:
        print(f"Error getting conversations: {str(e)}")
        return jsonify({"error": f"Failed to retrieve conversations: {str(e)}"}), 500

@app.route('/api/bots/<bot_id>/conversations/<conversation_id>', methods=['GET'])
@require_auth
def get_conversation(user_data, bot_id, conversation_id):
    """Get a specific conversation with all messages"""
    try:
        # Check if conversation exists
        conversation_file = os.path.join(CONVERSATIONS_FOLDER, f"{user_data['id']}_{bot_id}_{conversation_id}.json")
        if not os.path.exists(conversation_file):
            return jsonify({"error": "Conversation not found"}), 404
        
        # Read conversation
        with open(conversation_file, 'r') as f:
            conversation = json.load(f)
        
        return jsonify(conversation), 200
    except Exception as e:
        print(f"Error getting conversation: {str(e)}")
        return jsonify({"error": f"Failed to retrieve conversation: {str(e)}"}), 500

@app.route('/api/bots/<bot_id>/conversations/<conversation_id>', methods=['DELETE'])
@require_auth
def delete_conversation(user_data, bot_id, conversation_id):
    """Delete a specific conversation"""
    try:
        # Check if conversation exists
        conversation_file = os.path.join(CONVERSATIONS_FOLDER, f"{user_data['id']}_{bot_id}_{conversation_id}.json")
        if not os.path.exists(conversation_file):
            return jsonify({"error": "Conversation not found"}), 404
        
        # Delete the file
        os.remove(conversation_file)
        
        return jsonify({"message": "Conversation deleted successfully"}), 200
    except Exception as e:
        print(f"Error deleting conversation: {str(e)}")
        return jsonify({"error": f"Failed to delete conversation: {str(e)}"}), 500

@app.route('/api/bots/<bot_id>/conversations', methods=['DELETE'])
@require_auth
def delete_all_conversations(user_data, bot_id):
    """Delete all conversations for a specific bot"""
    try:
        # Find all conversation files for this user and bot
        conversation_files = [f for f in os.listdir(CONVERSATIONS_FOLDER) 
                             if f.startswith(f"{user_data['id']}_{bot_id}_") and f.endswith('.json')]
        
        # Delete each file
        for file_name in conversation_files:
            file_path = os.path.join(CONVERSATIONS_FOLDER, file_name)
            os.remove(file_path)
        
        return jsonify({"message": f"Deleted {len(conversation_files)} conversations successfully"}), 200
    except Exception as e:
        print(f"Error deleting all conversations: {str(e)}")
        return jsonify({"error": f"Failed to delete conversations: {str(e)}"}), 500

@app.route('/api/bots/<bot_id>/conversations/<conversation_id>/rename', methods=['POST'])
@require_auth
def rename_conversation(user_data, bot_id, conversation_id):
    """Rename a conversation"""
    try:
        data = request.json
        new_title = data.get('title')
        
        if not new_title:
            return jsonify({"error": "Title is required"}), 400
        
        # Check if conversation exists
        conversation_file = os.path.join(CONVERSATIONS_FOLDER, f"{user_data['id']}_{bot_id}_{conversation_id}.json")
        if not os.path.exists(conversation_file):
            return jsonify({"error": "Conversation not found"}), 404
        
        # Read conversation
        with open(conversation_file, 'r') as f:
            conversation = json.load(f)
        
        # Update title
        conversation["title"] = new_title
        
        # Save updated conversation
        with open(conversation_file, 'w') as f:
            json.dump(conversation, f)
        
        return jsonify({"message": "Conversation renamed successfully", "title": new_title}), 200
    except Exception as e:
        print(f"Error renaming conversation: {str(e)}")
        return jsonify({"error": f"Failed to rename conversation: {str(e)}"}), 500

# Add model chorus API endpoints
@app.route('/api/bots/<bot_id>/chorus', methods=['GET'])
@require_auth
def get_chorus_config(user_data, bot_id):
    # Check if bot exists and belongs to user
    bots_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "bots")
    user_bots_file = os.path.join(bots_dir, f"{user_data['id']}_bots.json")
    
    if not os.path.exists(user_bots_file):
        return jsonify({"error": "Bot not found"}), 404
        
    with open(user_bots_file, 'r') as f:
        bots = json.load(f)
        
    bot = None
    for b in bots:
        if b["id"] == bot_id:
            bot = b
            break
            
    if not bot:
        return jsonify({"error": "Bot not found"}), 404
    
    # Check for chorus_id in the bot
    chorus_id = bot.get("chorus_id", "")
    if not chorus_id:
        return jsonify({"error": "No chorus configuration found for this bot"}), 404
    
    # Get the user's chorus definitions
    chorus_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "chorus")
    os.makedirs(chorus_dir, exist_ok=True)
    
    user_choruses_file = os.path.join(chorus_dir, f"{user_data['id']}_choruses.json")
    if not os.path.exists(user_choruses_file):
        return jsonify({"error": "No chorus configurations found"}), 404
    
    # Load the user's choruses and find the one with the matching ID
    with open(user_choruses_file, 'r') as f:
        choruses = json.load(f)
    
    chorus = next((c for c in choruses if c["id"] == chorus_id), None)
    if not chorus:
        return jsonify({"error": f"Chorus configuration with ID {chorus_id} not found"}), 404
    
    return jsonify(chorus), 200

@app.route('/api/bots/<bot_id>/chorus', methods=['POST'])
@require_auth
def save_chorus_config(user_data, bot_id):
    # This endpoint is now simplified to just associate a chorus with a bot
    # For backward compatibility
    
    # Check if bot exists and belongs to user
    bots_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "bots")
    user_bots_file = os.path.join(bots_dir, f"{user_data['id']}_bots.json")
    
    if not os.path.exists(user_bots_file):
        return jsonify({"error": "Bot not found"}), 404
        
    with open(user_bots_file, 'r') as f:
        bots = json.load(f)
        
    bot = None
    bot_index = -1
    for i, b in enumerate(bots):
        if b["id"] == bot_id:
            bot = b
            bot_index = i
            break
            
    if not bot:
        return jsonify({"error": "Bot not found"}), 404
    
    # Get the chorus data
    data = request.json
    
    # First, create or update a chorus
    try:
        # Create chorus directory if it doesn't exist
        chorus_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "chorus")
        os.makedirs(chorus_dir, exist_ok=True)
        
        # Load or create user's chorus list
        user_choruses_file = os.path.join(chorus_dir, f"{user_data['id']}_choruses.json")
        
        if os.path.exists(user_choruses_file):
            with open(user_choruses_file, 'r') as f:
                choruses = json.load(f)
        else:
            choruses = []
        
        # Generate a unique ID for the chorus
        chorus_id = str(uuid.uuid4())
        
        # Create chorus data
        chorus_data = {
            "id": chorus_id,
            "name": data.get('name', 'Unnamed Chorus'),
            "description": data.get('description', ''),
            "response_models": data.get('response_models', []),
            "evaluator_models": data.get('evaluator_models', []),
            "created_at": datetime.datetime.now(UTC).isoformat(),
            "updated_at": datetime.datetime.now(UTC).isoformat(),
            "created_by": user_data['username']
        }
        
        # Add the chorus
        choruses.append(chorus_data)
        
        # Save the updated choruses
        with open(user_choruses_file, 'w') as f:
            json.dump(choruses, f)
            
        # Update the bot to use this chorus
        bots[bot_index]["chorus_id"] = chorus_id
        
        # Save the updated bot
        with open(user_bots_file, 'w') as f:
            json.dump(bots, f)
        
        return jsonify({
            "message": "Chorus created and assigned to bot successfully", 
            "config": chorus_data,
            "bot": bots[bot_index]
        }), 200
    except Exception as e:
        print(f"Error creating chorus: {str(e)}")
        return jsonify({"error": f"Error creating chorus: {str(e)}"}), 500

@app.route('/api/bots/<bot_id>/debug-flowchart', methods=['POST'])
@require_auth
def generate_debug_flowchart(user_data, bot_id):
    data = request.json
    message = data.get('message')
    use_model_chorus = data.get('use_model_chorus', False)
    
    if not message:
        return jsonify({"error": "Message is required"}), 400
        
    # Get bot info
    bots_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "bots")
    user_bots_file = os.path.join(bots_dir, f"{user_data['id']}_bots.json")
    
    if not os.path.exists(user_bots_file):
        return jsonify({"error": "Bot not found"}), 404
        
    with open(user_bots_file, 'r') as f:
        bots = json.load(f)
        
    bot = None
    for b in bots:
        if b["id"] == bot_id:
            bot = b
            break
            
    if not bot:
        return jsonify({"error": "Bot not found"}), 404
        
    # Get dataset info
    dataset_id = bot["dataset_id"]
    
    # Retrieve relevant documents from ChromaDB
    try:
        try:
            collection = chroma_client.get_collection(name=dataset_id, embedding_function=openai_ef)
        except Exception as coll_error:
            return jsonify({
                "error": f"Collection '{dataset_id}' not found."
            }), 404
                
        # Check if collection has any documents
        collection_count = collection.count()
        if collection_count == 0:
            return jsonify({
                "error": "No documents in collection"
            }), 404
                
        # Determine how many results to request based on collection size
        n_results = min(5, collection_count)
        
        results = collection.query(
            query_texts=[message],
            n_results=n_results
        )
        
        contexts = results["documents"][0]
        context_text = "\n\n".join(contexts)
        
        # Prepare prompt
        system_instruction = bot["system_instruction"]
        
        # Generate responses using model chorus if enabled
        if use_model_chorus:
            # Get responses from multiple models
            all_responses = []
            logs = ["Model chorus mode enabled"]
            
            # Function to generate OpenAI responses
            def get_openai_responses(count=3):
                openai_responses = []
                for i in range(count):
                    try:
                        response = openai.chat.completions.create(
                            model="gpt-4",
                            messages=[
                                {"role": "system", "content": system_instruction},
                                {"role": "user", "content": f"Context:\n{context_text}\n\nUser question: {message}"}
                            ],
                            temperature=0.5 + (i * 0.25),
                        )
                        response_text = response.choices[0].message.content
                        openai_responses.append({
                            "provider": "OpenAI",
                            "model": "GPT-4",
                            "temperature": 0.5 + (i * 0.25),
                            "response": response_text
                        })
                    except Exception as e:
                        logs.append(f"Error getting OpenAI response: {str(e)}")
                return openai_responses
                
            # Function to generate Anthropic responses
            def get_anthropic_responses(count=3):
                anthropic_responses = []
                for i in range(count):
                    try:
                        response = anthropic_client.messages.create(
                            model="claude-3-opus-20240229",
                            system=system_instruction,
                            messages=[
                                {"role": "user", "content": f"Context:\n{context_text}\n\nUser question: {message}"}
                            ],
                            temperature=0.5 + (i * 0.25),
                            max_tokens=1024
                        )
                        response_text = response.content[0].text
                        anthropic_responses.append({
                            "provider": "Anthropic",
                            "model": "Claude 3 Opus",
                            "temperature": 0.5 + (i * 0.25),
                            "response": response_text
                        })
                    except Exception as e:
                        logs.append(f"Error getting Anthropic response: {str(e)}")
                return anthropic_responses
                
            # Function to generate Groq responses
            def get_groq_responses(count=3):
                groq_responses = []
                for i in range(count):
                    try:
                        headers = {
                            "Authorization": f"Bearer {GROQ_API_KEY}",
                            "Content-Type": "application/json"
                        }
                        payload = {
                            "model": "llama3-70b-8192",
                            "messages": [
                                {"role": "system", "content": system_instruction},
                                {"role": "user", "content": f"Context:\n{context_text}\n\nUser question: {message}"}
                            ],
                            "temperature": 0.5 + (i * 0.25),
                            "max_tokens": 1024
                        }
                        response = requests.post(GROQ_API_URL, json=payload, headers=headers)
                        response_json = response.json()
                        response_text = response_json["choices"][0]["message"]["content"]
                        groq_responses.append({
                            "provider": "Groq",
                            "model": "LLaMA 3 70B",
                            "temperature": 0.5 + (i * 0.25),
                            "response": response_text
                        })
                    except Exception as e:
                        logs.append(f"Error getting Groq response: {str(e)}")
                return groq_responses
            
            # Gather responses from all providers
            openai_responses = get_openai_responses(3)
            anthropic_responses = get_anthropic_responses(3)
            groq_responses = get_groq_responses(3)
            
            # Combine all responses
            all_responses = openai_responses + anthropic_responses + groq_responses
            responses = [r["response"] for r in all_responses]
            
            # Generate Mermaid diagram code with provider information
            try:
                # Create a model chorus flowchart
                mermaid_code = "flowchart TD\n"
                
                # Define node IDs clearly to avoid conflicts
                mermaid_code += "    userQuery[\"User Query\"]\n"
                mermaid_code += "    contextRetrieval[\"Context Retrieval\"]\n"
                
                # Define provider groups
                mermaid_code += "    subgraph OpenAI\n"
                for i in range(len(openai_responses)):
                    mermaid_code += f"        openai{i+1}[\"OpenAI {i+1}\"]\n"
                mermaid_code += "    end\n"
                
                mermaid_code += "    subgraph Anthropic\n"
                for i in range(len(anthropic_responses)):
                    mermaid_code += f"        anthropic{i+1}[\"Claude {i+1}\"]\n"
                mermaid_code += "    end\n"
                
                mermaid_code += "    subgraph Groq\n"
                for i in range(len(groq_responses)):
                    mermaid_code += f"        groq{i+1}[\"LLaMA {i+1}\"]\n"
                mermaid_code += "    end\n"
                
                # Define evaluator node
                mermaid_code += "    evaluator[\"Voting Process\"]\n"
                
                # Define winner and final response
                mermaid_code += "    winner[\"Winning Response\"]\n"
                
                # Connect nodes
                mermaid_code += "    userQuery --> contextRetrieval\n"
                
                # Connect to each model
                for i in range(len(openai_responses)):
                    mermaid_code += f"    contextRetrieval --> openai{i+1}\n"
                    mermaid_code += f"    openai{i+1} --> evaluator\n"
                
                for i in range(len(anthropic_responses)):
                    mermaid_code += f"    contextRetrieval --> anthropic{i+1}\n"
                    mermaid_code += f"    anthropic{i+1} --> evaluator\n"
                
                for i in range(len(groq_responses)):
                    mermaid_code += f"    contextRetrieval --> groq{i+1}\n"
                    mermaid_code += f"    groq{i+1} --> evaluator\n"
                
                # Connect to winner
                mermaid_code += "    evaluator --> winner\n"
                
                # Add styling
                mermaid_code += "    classDef openai fill:#c1edc9,stroke:#333;\n"
                mermaid_code += "    classDef anthropic fill:#d4e1f5,stroke:#333;\n"
                mermaid_code += "    classDef groq fill:#f9d6c4,stroke:#333;\n"
                mermaid_code += "    classDef winner fill:#f9f,stroke:#333,stroke-width:4px;\n"
                
                # Apply classes
                for i in range(len(openai_responses)):
                    mermaid_code += f"    class openai{i+1} openai;\n"
                
                for i in range(len(anthropic_responses)):
                    mermaid_code += f"    class anthropic{i+1} anthropic;\n"
                
                for i in range(len(groq_responses)):
                    mermaid_code += f"    class groq{i+1} groq;\n"
                
                mermaid_code += "    class winner winner;\n"
                
                return jsonify({
                    "mermaid_code": mermaid_code,
                    "data": {
                        "responses": responses,
                        "all_responses": all_responses,
                        "logs": logs,
                        "contexts": contexts
                    }
                }), 200
            except Exception as mermaid_error:
                print(f"Error generating chorus Mermaid code: {str(mermaid_error)}")
                # Provide a simplified fallback diagram
                mermaid_code = """flowchart TD
    A[User Query] --> B[Multi-Model Processing]
    B --> C[Response]
    style C fill:#f9f,stroke:#333,stroke-width:4px"""
                
                return jsonify({
                    "mermaid_code": mermaid_code,
                    "data": {
                        "responses": responses,
                        "all_responses": all_responses,
                        "logs": logs,
                        "contexts": contexts,
                        "error": str(mermaid_error)
                    }
                }), 200
        else:
            # In debug mode - generate 5 responses and let them vote
            responses = []
            anonymized_responses = []
            response_metadata = []
            logs = ["Debug mode enabled - generating 5 different responses"]
            
            for i in range(5):
                logs.append(f"Generating response {i+1}/5...")
                response = openai.chat.completions.create(
                    model="gpt-4",
                    messages=[
                        {"role": "system", "content": system_instruction},
                        {"role": "user", "content": f"Context:\n{contexts}\n\nUser question: {message}"}
                    ],
                    temperature=0.7 + (i * 0.1),  # Vary temperature to get different responses
                )
                response_text = response.choices[0].message.content
                responses.append(response_text)
                anonymized_responses.append(response_text)
                response_metadata.append({
                    "index": i+1,
                    "temperature": 0.7 + (i * 0.1)
                })
                logs.append(f"Response {i+1}: {response_text[:100]}...")
                
            # Have the responses vote on each other
            logs.append("Having the responses evaluate each other...")
            votes = [0, 0, 0, 0, 0]
            vote_details = []
            
            for i in range(5):
                voting_prompt = f"""You are an expert evaluator. You need to rank the following responses to the question: "{message}"
                
Here are the 5 candidate responses:

Response 1:
{anonymized_responses[0]}

Response 2:
{anonymized_responses[1]}

Response 3:
{anonymized_responses[2]}

Response 4:
{anonymized_responses[3]}

Response 5:
{anonymized_responses[4]}

Which response provides the most accurate, helpful, and relevant answer? Return ONLY the number (1-5) of the best response.
Do not reveal any bias or preference based on writing style or approach - evaluate solely on answer quality, accuracy and helpfulness.
"""
                voting_response = openai.chat.completions.create(
                    model="gpt-4",
                    messages=[{"role": "user", "content": voting_prompt}],
                    temperature=0.2,
                )
                
                vote_text = voting_response.choices[0].message.content
                
                try:
                    # More robust vote extraction 
                    vote_number = None
                    vote_text_lower = vote_text.lower().strip()
                    
                    # Try regex first
                    match = re.search(r'\b([1-5])\b', vote_text_lower)
                    if match:
                        vote_number = int(match.group(1))
                    # If no match, look for number words
                    elif "one" in vote_text_lower or "first" in vote_text_lower:
                        vote_number = 1
                    elif "two" in vote_text_lower or "second" in vote_text_lower:
                        vote_number = 2
                    elif "three" in vote_text_lower or "third" in vote_text_lower:
                        vote_number = 3
                    elif "four" in vote_text_lower or "fourth" in vote_text_lower:
                        vote_number = 4
                    elif "five" in vote_text_lower or "fifth" in vote_text_lower:
                        vote_number = 5
                    # Fall back to the old method
                    else:
                        for j in range(1, 6):
                            if str(j) in vote_text:
                                vote_number = j
                                break
                    
                    # Record the vote
                    if vote_number and 1 <= vote_number <= 5:
                        votes[vote_number-1] += 1
                        logs.append(f"Evaluator {i+1} voted for Response {vote_number}")
                    else:
                        logs.append(f"Evaluator {i+1} vote unclear: {vote_text}")
                    
                    vote_details.append({
                        "evaluator": i+1,
                        "vote": vote_number,
                        "raw_text": vote_text
                    })
                except Exception as e:
                    logs.append(f"Couldn't parse vote from: {vote_text}. Error: {str(e)}")
                    vote_details.append({
                        "evaluator": i+1,
                        "vote": None,
                        "raw_text": vote_text
                    })
            
            # Log the final results
            logs.append(f"Voting results: {votes}")
            
            # Find the winning response
            max_votes = max(votes)
            winning_indices = [i for i, v in enumerate(votes) if v == max_votes]
            winning_index = winning_indices[0]  # Take the first one in case of a tie
            
            logs.append(f"Response {winning_index + 1} wins with {max_votes} votes")
            
            # Generate Mermaid diagram code - with safer, more reliable text escaping
            try:
                # Create a simpler, more reliable flowchart
                mermaid_code = "flowchart TD\n"
                
                # Define node IDs clearly to avoid conflicts
                mermaid_code += "    userQuery[\"User Query\"]\n"
                mermaid_code += "    contextRetrieval[\"Context Retrieval\"]\n"
                
                # Define response nodes
                for i in range(5):
                    mermaid_code += f"    response{i+1}[\"Response {i+1}\"]\n"
                
                # Define voting process node
                mermaid_code += "    votingProcess[\"Voting Process\"]\n"
                
                # Define evaluator nodes
                for i in range(5):
                    mermaid_code += f"    evaluator{i+1}[\"Evaluator {i+1}\"]\n"
                
                # Define results and final response nodes
                mermaid_code += "    results[\"Results\"]\n"
                mermaid_code += "    finalResponse[\"Final Response\"]\n"
                
                # Connect nodes
                mermaid_code += "    userQuery --> contextRetrieval\n"
                
                # Connect context to responses
                for i in range(5):
                    mermaid_code += f"    contextRetrieval --> response{i+1}\n"
                
                # Connect responses to voting process - ensure ALL responses connect to voting process
                for i in range(5):
                    mermaid_code += f"    response{i+1} --> votingProcess\n"
                
                # Connect voting process to evaluators
                for i in range(5):
                    mermaid_code += f"    votingProcess --> evaluator{i+1}\n"
                
                # Connect evaluators to results
                for i in range(5):
                    mermaid_code += f"    evaluator{i+1} --> results\n"
                
                # Connect results to final response
                mermaid_code += "    results --> finalResponse\n"
                
                # Label nodes with actual content (safely escaped)
                message_safe = message.replace('"', "'").replace('\n', ' ')[:30]
                mermaid_code += f"    userQuery[\"User Query:<br/>{message_safe}...\"]\n"
                mermaid_code += f"    contextRetrieval[\"Context Retrieval:<br/>{len(contexts)} chunks\"]\n"
                
                # Add response content
                for i in range(5):
                    resp_text = responses[i].replace('"', "'").replace('\n', ' ')[:30]
                    temp = 0.7 + (i * 0.1)
                    mermaid_code += f"    response{i+1}[\"Response {i+1} (Temp: {temp:.1f}):<br/>{resp_text}...\"]\n"
                
                # Add vote counts to results
                mermaid_code += f"    results[\"Results: {votes}<br/>Winner: Response {winning_index + 1} ({max_votes} votes)\"]\n"
                
                # Add winner content to final response
                final_text = responses[winning_index].replace('"', "'").replace('\n', ' ')[:30]
                mermaid_code += f"    finalResponse[\"Final Response:<br/>{final_text}...\"]\n"
                
                # Style nodes
                mermaid_code += "    class userQuery query;\n"
                mermaid_code += "    class contextRetrieval context;\n"
                mermaid_code += "    class response1,response2,response3,response4,response5 response;\n"
                mermaid_code += "    class evaluator1,evaluator2,evaluator3,evaluator4,evaluator5 evaluator;\n"
                mermaid_code += "    class results results;\n"
                mermaid_code += "    class finalResponse final;\n"
                mermaid_code += "    class votingProcess voting;\n"
                
                # Highlight winning response
                mermaid_code += f"    class response{winning_index + 1} winner;\n"
                
                # Add special connections for votes - only from evaluators to responses
                for i, vote_info in enumerate(vote_details):
                    vote = vote_info.get("vote")
                    if vote is not None:
                        # Use dashed lines for vote connections
                        mermaid_code += f"    evaluator{i+1} -.-> response{vote}\n"
                
                # Add styling
                mermaid_code += "    classDef query fill:#d4f1c5;\n"
                mermaid_code += "    classDef context fill:#fff4b8;\n"
                mermaid_code += "    classDef response fill:#ffcfcf;\n"
                mermaid_code += "    classDef evaluator fill:#ffd7b5;\n"
                mermaid_code += "    classDef results fill:#ffe066;\n"
                mermaid_code += "    classDef final fill:#b3e6cc;\n"
                mermaid_code += "    classDef voting fill:#ffaaaa;\n"
                mermaid_code += "    classDef winner fill:#ffcfcf,stroke:#ff0000,stroke-width:2px;\n"
            except Exception as mermaid_error:
                print(f"Error generating Mermaid code: {str(mermaid_error)}")
                # Provide a simplified fallback diagram
                mermaid_code = """flowchart TD
    A[User Query] --> B[Processing]
    B --> C[Response]
    style C fill:#f9f,stroke:#333,stroke-width:4px"""
            
            return jsonify({
                "mermaid_code": mermaid_code,
                "data": {
                    "responses": responses,
                    "anonymized_responses": anonymized_responses,
                    "response_metadata": response_metadata,
                    "votes": votes,
                    "logs": logs,
                    "contexts": contexts,
                    "vote_details": vote_details
                }
            }), 200
                
    except Exception as e:
        print(f"Error in generate_debug_flowchart: {str(e)}")
        return jsonify({"error": "Failed to generate flowchart", "details": str(e)}), 500

@app.route('/api/bots/<bot_id>/set-chorus', methods=['POST'])
@require_auth
def set_bot_chorus(user_data, bot_id):
    # Check if bot exists and belongs to user
    bots_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "bots")
    user_bots_file = os.path.join(bots_dir, f"{user_data['id']}_bots.json")
    
    if not os.path.exists(user_bots_file):
        return jsonify({"error": "Bot not found"}), 404
        
    with open(user_bots_file, 'r') as f:
        bots = json.load(f)
        
    bot = None
    bot_index = -1
    for i, b in enumerate(bots):
        if b["id"] == bot_id:
            bot = b
            bot_index = i
            break
            
    if not bot:
        return jsonify({"error": "Bot not found"}), 404
    
    # Get the chorus ID from the request
    data = request.json
    chorus_id = data.get('chorus_id', '')
    
    # Empty string means unassign any chorus
    if chorus_id:
        # Verify that the chorus exists if an ID was provided
        chorus_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "chorus")
        user_choruses_file = os.path.join(chorus_dir, f"{user_data['id']}_choruses.json")
        
        if os.path.exists(user_choruses_file):
            with open(user_choruses_file, 'r') as f:
                choruses = json.load(f)
                
            # Find the chorus by ID
            chorus = next((c for c in choruses if c["id"] == chorus_id), None)
            if not chorus:
                return jsonify({"error": f"Chorus with ID {chorus_id} not found"}), 404
    
    # Update the bot with the new chorus ID
    bots[bot_index]["chorus_id"] = chorus_id
    
    # Save the updated bot
    with open(user_bots_file, 'w') as f:
        json.dump(bots, f)
    
    return jsonify({
        "message": chorus_id and "Chorus assigned to bot successfully" or "Chorus unassigned from bot successfully", 
        "bot": bots[bot_index]
    }), 200

@app.route('/api/choruses', methods=['GET'])
@require_auth
def list_choruses(user_data):
    # Get all chorus configurations
    chorus_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "chorus")
    os.makedirs(chorus_dir, exist_ok=True)
    
    # First check if there's a user-specific choruses file
    user_choruses_file = os.path.join(chorus_dir, f"{user_data['id']}_choruses.json")
    
    # If the user choruses file doesn't exist yet, create it
    if not os.path.exists(user_choruses_file):
        with open(user_choruses_file, 'w') as f:
            json.dump([], f)
    
    # Load user's chorus definitions
    with open(user_choruses_file, 'r') as f:
        choruses = json.load(f)
    
    return jsonify(choruses), 200

@app.route('/api/choruses', methods=['POST'])
@require_auth
def create_chorus(user_data):
    # Get chorus data from request
    data = request.json
    
    # Basic validation
    if not data.get('name'):
        return jsonify({"error": "Chorus name is required"}), 400
        
    if not data.get('response_models') or not isinstance(data.get('response_models'), list) or len(data.get('response_models')) == 0:
        return jsonify({"error": "At least one response model is required"}), 400
        
    if not data.get('evaluator_models') or not isinstance(data.get('evaluator_models'), list) or len(data.get('evaluator_models')) == 0:
        return jsonify({"error": "At least one evaluator model is required"}), 400
    
    # Create chorus directory if it doesn't exist
    chorus_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "chorus")
    os.makedirs(chorus_dir, exist_ok=True)
    
    # Load user's chorus definitions
    user_choruses_file = os.path.join(chorus_dir, f"{user_data['id']}_choruses.json")
    if os.path.exists(user_choruses_file):
        with open(user_choruses_file, 'r') as f:
            choruses = json.load(f)
    else:
        choruses = []
    
    # Generate a unique ID for the chorus
    chorus_id = str(uuid.uuid4())
    
    # Add the new chorus with metadata
    chorus = {
        "id": chorus_id,
        "name": data.get('name'),
        "description": data.get('description', ''),
        "created_at": datetime.datetime.now(UTC).isoformat(),
        "updated_at": datetime.datetime.now(UTC).isoformat(),
        "response_models": data.get('response_models', []),
        "evaluator_models": data.get('evaluator_models', []),
        "created_by": user_data['username']
    }
    
    choruses.append(chorus)
    
    # Save updated chorus list
    with open(user_choruses_file, 'w') as f:
        json.dump(choruses, f)
    
    return jsonify(chorus), 201

@app.route('/api/choruses/<chorus_id>', methods=['GET'])
@require_auth
def get_chorus(user_data, chorus_id):
    # Get chorus directory and user choruses file
    chorus_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "chorus")
    user_choruses_file = os.path.join(chorus_dir, f"{user_data['id']}_choruses.json")
    
    if not os.path.exists(user_choruses_file):
        return jsonify({"error": "Chorus not found"}), 404
    
    # Load user's chorus definitions
    with open(user_choruses_file, 'r') as f:
        choruses = json.load(f)
    
    # Find the requested chorus
    chorus = next((c for c in choruses if c["id"] == chorus_id), None)
    if not chorus:
        return jsonify({"error": "Chorus not found"}), 404
    
    return jsonify(chorus), 200

@app.route('/api/choruses/<chorus_id>', methods=['PUT'])
@require_auth
def update_chorus(user_data, chorus_id):
    # Get chorus data from request
    data = request.json
    
    # Basic validation
    if not data.get('name'):
        return jsonify({"error": "Chorus name is required"}), 400
        
    if not data.get('response_models') or not isinstance(data.get('response_models'), list) or len(data.get('response_models')) == 0:
        return jsonify({"error": "At least one response model is required"}), 400
        
    if not data.get('evaluator_models') or not isinstance(data.get('evaluator_models'), list) or len(data.get('evaluator_models')) == 0:
        return jsonify({"error": "At least one evaluator model is required"}), 400
    
    # Get chorus directory and user choruses file
    chorus_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "chorus")
    user_choruses_file = os.path.join(chorus_dir, f"{user_data['id']}_choruses.json")
    
    if not os.path.exists(user_choruses_file):
        return jsonify({"error": "Chorus not found"}), 404
    
    # Load user's chorus definitions
    with open(user_choruses_file, 'r') as f:
        choruses = json.load(f)
    
    # Find and update the chorus
    chorus_index = next((i for i, c in enumerate(choruses) if c["id"] == chorus_id), None)
    if chorus_index is None:
        return jsonify({"error": "Chorus not found"}), 404
    
    # Update the chorus with new data, preserving the ID and creation date
    chorus = choruses[chorus_index]
    updated_chorus = {
        "id": chorus["id"],
        "name": data.get('name'),
        "description": data.get('description', ''),
        "created_at": chorus["created_at"],
        "updated_at": datetime.datetime.now(UTC).isoformat(),
        "response_models": data.get('response_models', []),
        "evaluator_models": data.get('evaluator_models', []),
        "created_by": chorus.get("created_by", user_data['username'])
    }
    
    choruses[chorus_index] = updated_chorus
    
    # Save updated chorus list
    with open(user_choruses_file, 'w') as f:
        json.dump(choruses, f)
    
    return jsonify(updated_chorus), 200

@app.route('/api/choruses/<chorus_id>', methods=['DELETE'])
@require_auth
def delete_chorus(user_data, chorus_id):
    # Get chorus directory and user choruses file
    chorus_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "chorus")
    user_choruses_file = os.path.join(chorus_dir, f"{user_data['id']}_choruses.json")
    
    if not os.path.exists(user_choruses_file):
        return jsonify({"error": "Chorus not found"}), 404
    
    # Load user's chorus definitions
    with open(user_choruses_file, 'r') as f:
        choruses = json.load(f)
    
    # Find and remove the chorus
    original_length = len(choruses)
    choruses = [c for c in choruses if c["id"] != chorus_id]
    
    if len(choruses) == original_length:
        return jsonify({"error": "Chorus not found"}), 404
    
    # Save updated chorus list
    with open(user_choruses_file, 'w') as f:
        json.dump(choruses, f)
    
    # Also check if any bots are using this chorus and update them
    bots_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "bots")
    user_bots_file = os.path.join(bots_dir, f"{user_data['id']}_bots.json")
    
    if os.path.exists(user_bots_file):
        with open(user_bots_file, 'r') as f:
            bots = json.load(f)
        
        # Update any bots using this chorus
        updated = False
        for bot in bots:
            if bot.get("chorus_id") == chorus_id:
                bot["chorus_id"] = ""
                updated = True
        
        # Save updated bots if needed
        if updated:
            with open(user_bots_file, 'w') as f:
                json.dump(bots, f)
    
    return jsonify({"message": "Chorus deleted successfully"}), 200

# Bots routes
@app.route('/api/bots', methods=['GET'])
@require_auth
def get_bots(user_data):
    # Get bots directory and user bots file
    bots_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "bots")
    os.makedirs(bots_dir, exist_ok=True)
    user_bots_file = os.path.join(bots_dir, f"{user_data['id']}_bots.json")
    
    # If user doesn't have any bots yet, return empty list
    if not os.path.exists(user_bots_file):
        return jsonify([]), 200
        
    # Return user's bots
    with open(user_bots_file, 'r') as f:
        bots = json.load(f)
    
    return jsonify(bots), 200

@app.route('/api/bots', methods=['POST'])
@require_auth
def create_bot(user_data):
    # Get bot data from request
    data = request.json
    
    if not data or not data.get('name'):
        return jsonify({"error": "Bot name is required"}), 400
    
    # Process dataset IDs - can be added later or supplied in the request
    dataset_ids = []
    if data.get('dataset_id'):  # Support single dataset_id in request for backward compatibility
        dataset_ids.append(data.get('dataset_id'))
    if data.get('dataset_ids') and isinstance(data.get('dataset_ids'), list):
        dataset_ids.extend(data.get('dataset_ids'))
    
    # Get chorus_id if provided
    chorus_id = data.get('chorus_id', '')
    
    # Verify that the chorus exists if one was provided
    if chorus_id:
        chorus_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "chorus")
        user_choruses_file = os.path.join(chorus_dir, f"{user_data['id']}_choruses.json")
        
        if os.path.exists(user_choruses_file):
            with open(user_choruses_file, 'r') as f:
                choruses = json.load(f)
                
            # Find the chorus by ID
            chorus = next((c for c in choruses if c["id"] == chorus_id), None)
            if not chorus:
                return jsonify({"error": f"Chorus with ID {chorus_id} not found"}), 400
    
    # Create a new bot
    new_bot = {
        "id": str(uuid.uuid4()),
        "name": data.get('name'),
        "description": data.get('description', ''),
        "dataset_ids": dataset_ids,
        "chorus_id": chorus_id,  # Set the chorus ID
        "prompt_template": data.get('prompt_template', ''),
        "system_instruction": data.get('system_instruction', 'You are a helpful AI assistant. Answer questions based on the provided context.'),
        "created_at": datetime.datetime.now(UTC).isoformat()
    }
    
    # Save the bot
    bots_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "bots")
    os.makedirs(bots_dir, exist_ok=True)
    
    # Check if user already has bots
    user_bots_file = os.path.join(bots_dir, f"{user_data['id']}_bots.json")
    
    if os.path.exists(user_bots_file):
        with open(user_bots_file, 'r') as f:
            bots = json.load(f)
        bots.append(new_bot)
    else:
        bots = [new_bot]
    
    with open(user_bots_file, 'w') as f:
        json.dump(bots, f)
    
    return jsonify(new_bot), 201

@app.route('/api/bots/<bot_id>', methods=['DELETE'])
@require_auth
def delete_bot(user_data, bot_id):
    # Get bots directory and user bots file
    bots_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "bots")
    user_bots_file = os.path.join(bots_dir, f"{user_data['id']}_bots.json")
    
    if not os.path.exists(user_bots_file):
        return jsonify({"error": "Bot not found"}), 404
        
    with open(user_bots_file, 'r') as f:
        bots = json.load(f)
    
    # Find and remove the bot
    bot_found = False
    for i, bot in enumerate(bots):
        if bot["id"] == bot_id:
            bots.pop(i)
            bot_found = True
            break
    
    if not bot_found:
        return jsonify({"error": "Bot not found"}), 404
    
    # Save updated bots list
    with open(user_bots_file, 'w') as f:
        json.dump(bots, f)
    
    return jsonify({"message": "Bot deleted successfully"}), 200

# Add new endpoints for managing datasets in a bot
@app.route('/api/bots/<bot_id>/datasets', methods=['GET'])
@require_auth
def get_bot_datasets(user_data, bot_id):
    # Get bot info
    bots_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "bots")
    user_bots_file = os.path.join(bots_dir, f"{user_data['id']}_bots.json")
    
    if not os.path.exists(user_bots_file):
        return jsonify({"error": "Bot not found"}), 404
        
    with open(user_bots_file, 'r') as f:
        bots = json.load(f)
        
    bot = None
    for b in bots:
        if b["id"] == bot_id:
            bot = b
            break
            
    if not bot:
        return jsonify({"error": "Bot not found"}), 404
    
    # Get dataset info for each dataset ID
    datasets_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "datasets")
    user_datasets_file = os.path.join(datasets_dir, f"{user_data['id']}_datasets.json")
    
    dataset_details = []
    
    # Get dataset IDs from bot
    dataset_ids = bot.get("dataset_ids", [])
    
    if os.path.exists(user_datasets_file):
        try:
            with open(user_datasets_file, 'r') as f:
                all_datasets = json.load(f)
                
            for dataset_id in dataset_ids:
                dataset_found = False
                for d in all_datasets:
                    if d["id"] == dataset_id:
                        dataset_details.append(d)
                        dataset_found = True
                        break
                
                if not dataset_found:
                    # Add placeholder for missing dataset
                    dataset_details.append({
                        "id": dataset_id,
                        "name": "Unknown dataset",
                        "missing": True
                    })
        except Exception as e:
            print(f"Error loading datasets: {str(e)}")
    
    return jsonify({
        "bot": bot,
        "datasets": dataset_details
    }), 200

@app.route('/api/bots/<bot_id>/datasets', methods=['POST'])
@require_auth
def add_dataset_to_bot(user_data, bot_id):
    data = request.json
    dataset_id = data.get('dataset_id')
    
    if not dataset_id:
        return jsonify({"error": "Dataset ID is required"}), 400
    
    # Check if dataset exists
    datasets_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "datasets")
    user_datasets_file = os.path.join(datasets_dir, f"{user_data['id']}_datasets.json")
    
    dataset_exists = False
    dataset_name = "Unknown dataset"
    if os.path.exists(user_datasets_file):
        with open(user_datasets_file, 'r') as f:
            datasets = json.load(f)
        
        for d in datasets:
            if d["id"] == dataset_id:
                dataset_exists = True
                dataset_name = d["name"]
                break
    
    if not dataset_exists:
        return jsonify({"error": "Dataset not found"}), 404
    
    # Get bot info
    bots_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "bots")
    user_bots_file = os.path.join(bots_dir, f"{user_data['id']}_bots.json")
    
    if not os.path.exists(user_bots_file):
        return jsonify({"error": "Bot not found"}), 404
        
    with open(user_bots_file, 'r') as f:
        bots = json.load(f)
        
    bot_index = None
    for i, b in enumerate(bots):
        if b["id"] == bot_id:
            bot_index = i
            break
            
    if bot_index is None:
        return jsonify({"error": "Bot not found"}), 404
    
    # Update bot with the new dataset
    bot = bots[bot_index]
    
    # Set up dataset_ids if it doesn't exist
    if "dataset_ids" not in bot:
        bot["dataset_ids"] = []
        
    # Add dataset if not already added
    if dataset_id not in bot["dataset_ids"]:
        bot["dataset_ids"].append(dataset_id)
        
    # Save updated bot
    bots[bot_index] = bot
    with open(user_bots_file, 'w') as f:
        json.dump(bots, f)
    
    return jsonify({
        "message": f"Dataset '{dataset_name}' added to bot successfully",
        "bot": bot
    }), 200

@app.route('/api/bots/<bot_id>/datasets/<dataset_id>', methods=['DELETE'])
@require_auth
def remove_dataset_from_bot(user_data, bot_id, dataset_id):
    # Get bot info
    bots_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "bots")
    user_bots_file = os.path.join(bots_dir, f"{user_data['id']}_bots.json")
    
    if not os.path.exists(user_bots_file):
        return jsonify({"error": "Bot not found"}), 404
        
    with open(user_bots_file, 'r') as f:
        bots = json.load(f)
        
    bot_index = None
    for i, b in enumerate(bots):
        if b["id"] == bot_id:
            bot_index = i
            break
            
    if bot_index is None:
        return jsonify({"error": "Bot not found"}), 404
    
    # Update bot to remove the dataset
    bot = bots[bot_index]
    
    # Remove dataset from the dataset_ids array
    dataset_removed = False
    if "dataset_ids" in bot and dataset_id in bot["dataset_ids"]:
        bot["dataset_ids"].remove(dataset_id)
        dataset_removed = True
    
    if not dataset_removed:
        return jsonify({"error": "Dataset not found in bot"}), 404
    
    # Save updated bot
    bots[bot_index] = bot
    with open(user_bots_file, 'w') as f:
        json.dump(bots, f)
    
    return jsonify({
        "message": "Dataset removed from bot successfully",
        "bot": bot
    }), 200

# Image generation route
@app.route('/api/images/generate', methods=['POST'])
@require_auth
def generate_image(user_data):
    try:
        # Get request data
        data = request.json
        prompt = data.get('prompt')
        
        if not prompt:
            return jsonify({"error": "Prompt is required"}), 400
            
        # Get additional parameters with defaults for gpt-image-1
        model = "gpt-image-1"  # Always use gpt-image-1
        size = data.get('size', '1024x1024')  # Default square
        quality = data.get('quality', 'auto')  # 'low', 'medium', 'high', or 'auto'
        n = data.get('n', 1)  # Number of images to generate
        output_format = data.get('output_format', 'png')  # 'png', 'jpeg', or 'webp'
        moderation = data.get('moderation', 'auto')  # 'auto' or 'low'
        
        # Build the API request params with only supported parameters
        generation_params = {
            "model": model,
            "prompt": prompt,
            "n": n,
            "size": size,
            "quality": quality,
            "moderation": moderation
        }
        
        # Generate image with OpenAI
        response = openai.images.generate(**generation_params)
        
        results = []
        # Process all generated images
        for i, image_obj in enumerate(response.data):
            # Get the base64 content from OpenAI response
            if hasattr(image_obj, 'b64_json') and image_obj.b64_json:
                # Decode base64 content
                image_content = base64.b64decode(image_obj.b64_json)
                image_url = None
            elif hasattr(image_obj, 'url') and image_obj.url:
                # For backward compatibility if URL is provided
                image_url = image_obj.url
                
                # Download the image
                image_response = requests.get(image_url)
                if image_response.status_code != 200:
                    return jsonify({"error": f"Failed to download generated image {i+1}"}), 500
                    
                image_content = image_response.content
            else:
                return jsonify({"error": f"No image data found in response for image {i+1}"}), 500
                
            # Create a unique filename
            ext = output_format if output_format else "png"  # Default to png
            filename = f"generated_{str(uuid.uuid4())}.{ext}"
            filepath = os.path.join(IMAGE_FOLDER, filename)
            
            # Save the image
            with open(filepath, 'wb') as f:
                f.write(image_content)
                
            # Create the URL for accessing the image
            api_image_url = f"/api/images/{filename}"
            
            # Add to results
            results.append({
                "image_url": api_image_url,
                "original_url": image_url,
                "params": {
                    "model": model,
                    "size": size,
                    "quality": quality,
                    "format": output_format
                }
            })
        
        return jsonify({
            "success": True,
            "images": results
        })
        
    except Exception as e:
        print(f"Error generating image: {str(e)}")
        return jsonify({"error": f"Image generation failed: {str(e)}"}), 500

# Add a new endpoint for enhancing prompts with GPT-4o
@app.route('/api/images/enhance-prompt', methods=['POST'])
@require_auth
def enhance_prompt(user_data):
    try:
        print("\n=== Starting Prompt Enhancement ===")
        # Get request data
        data = request.json
        original_prompt = data.get('prompt')
        print(f"Received prompt: {original_prompt}")
        
        if not original_prompt:
            print("Error: No prompt provided")
            return jsonify({"error": "Prompt is required"}), 400
        
        # Craft system message for GPT-4o to enhance the prompt
        system_message = """You are an expert at creating detailed image prompts for AI image generation. 
Your task is to enhance the user's prompt by adding more descriptive elements, artistic style, lighting, mood, and details.
Keep the original intent and subject matter, but make it much more detailed and visually compelling.
Respond with ONLY the enhanced prompt text, nothing else. No explanations or additional text."""

        print("Calling GPT-4o for enhancement...")
        # Call GPT-4o to enhance the prompt
        response = openai.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": system_message},
                {"role": "user", "content": f"Enhance this image prompt: {original_prompt}"}
            ],
            temperature=0.7,
            max_tokens=500
        )
        
        enhanced_prompt = response.choices[0].message.content.strip()
        print(f"Enhanced prompt: {enhanced_prompt}")
        
        result = {
            "success": True,
            "original_prompt": original_prompt,
            "enhanced_prompt": enhanced_prompt
        }
        print("=== Prompt Enhancement Complete ===\n")
        return jsonify(result), 200
        
    except Exception as e:
        print(f"Error enhancing prompt: {str(e)}")
        print(f"Full error details: {e}")
        print("=== Prompt Enhancement Failed ===\n")
        return jsonify({"error": f"Prompt enhancement failed: {str(e)}"}), 500
        
# Route to serve generated images
@app.route('/api/images/<filename>', methods=['GET'])
def get_image(filename):
    """Serve an uploaded image file"""
    # Check if download parameter is provided
    download = request.args.get('download', 'false').lower() == 'true'
    
    if download:
        # If download is requested, set Content-Disposition header
        return send_from_directory(
            IMAGE_FOLDER, 
            filename, 
            as_attachment=True,
            download_name=filename
        )
    else:
        # Regular image view
        return send_from_directory(IMAGE_FOLDER, filename)

# Add this new endpoint after the existing chat_with_bot function
@app.route('/api/bots/<bot_id>/chat-with-image', methods=['POST'])
@require_auth
def chat_with_image(user_data, bot_id):
    # Check if the request is multipart/form-data or JSON
    if request.content_type and request.content_type.startswith('multipart/form-data'):
        # Handle multipart/form-data request
        message = request.form.get('message', '')
        debug_mode = request.form.get('debug_mode', 'false').lower() == 'true'
        conversation_id = request.form.get('conversation_id', '')
        
        # Get the image file
        if 'image' not in request.files:
            return jsonify({"error": "Image file is required"}), 400
            
        image_file = request.files['image']
        if image_file.filename == '':
            return jsonify({"error": "No image selected"}), 400
        
        # Create a new conversation if no conversation_id provided
        if not conversation_id:
            conversation_id = str(uuid.uuid4())
        
        # Process the image file
        try:
            # Save the image to a temporary file
            filename = f"chat_image_{str(uuid.uuid4())}{os.path.splitext(image_file.filename)[1]}"
            image_path = os.path.join(IMAGE_FOLDER, filename)
            image_file.save(image_path)
            
            # Get the file extension
            file_ext = os.path.splitext(image_file.filename)[1][1:]  # Remove the dot
            if not file_ext:
                file_ext = 'png'  # Default to PNG if no extension
        
            # Continue with image processing
        except Exception as img_error:
            print(f"Error processing image: {str(img_error)}")
            return jsonify({"error": f"Failed to process image: {str(img_error)}"}), 400
    else:
        # Handle JSON request
        data = request.json
        if not data:
            return jsonify({"error": "Invalid request format"}), 400
            
        message = data.get('message', '')
        image_data = data.get('image_data', '')
        debug_mode = data.get('debug_mode', False)
        conversation_id = data.get('conversation_id', '')
        
        if not image_data:
            return jsonify({"error": "Image data is required"}), 400
        
        # Create a new conversation if no conversation_id provided
        if not conversation_id:
            conversation_id = str(uuid.uuid4())
        
        # Process the base64 image
        try:
            # Extract the base64 content and file extension
            if ';base64,' in image_data:
                header, encoded = image_data.split(';base64,')
                file_ext = header.split('/')[-1]
            else:
                encoded = image_data
                file_ext = 'png'  # Default to PNG if not specified
            
            # Decode the base64 data
            decoded_image = base64.b64decode(encoded)
            
            # Save to a temporary file
            filename = f"chat_image_{str(uuid.uuid4())}.{file_ext}"
            image_path = os.path.join(IMAGE_FOLDER, filename)
            
            with open(image_path, 'wb') as f:
                f.write(decoded_image)
        except Exception as img_error:
            print(f"Error processing image: {str(img_error)}")
            return jsonify({"error": f"Failed to process image: {str(img_error)}"}), 400
    
    # Resize image if needed
    image_path = resize_image(image_path)
    
    # Add message text if not provided
    if not message:
        message = "[Image uploaded]"
    
    # Get bot info
    bots_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "bots")
    user_bots_file = os.path.join(bots_dir, f"{user_data['id']}_bots.json")
    
    if not os.path.exists(user_bots_file):
        return jsonify({"error": "Bot not found"}), 404
        
    with open(user_bots_file, 'r') as f:
        bots = json.load(f)
        
    bot = None
    for b in bots:
        if b["id"] == bot_id:
            bot = b
            break
            
    if not bot:
        return jsonify({"error": "Bot not found"}), 404
    
    # Create message object
    user_message = {
        "id": str(uuid.uuid4()),
        "role": "user",
        "content": message,
        "timestamp": datetime.datetime.now(UTC).isoformat(),
        "has_image": True,
        "image_path": image_path
    }
    
    # Add message to conversation history
    conversation_file = os.path.join(CONVERSATIONS_FOLDER, f"{user_data['id']}_{bot_id}_{conversation_id}.json")
    
    if os.path.exists(conversation_file):
        # Load existing conversation
        with open(conversation_file, 'r') as f:
            conversation = json.load(f)
    else:
        # Create new conversation
        conversation = {
            "id": conversation_id,
            "bot_id": bot_id,
            "user_id": user_data['id'],
            "title": message[:40] + "..." if len(message) > 40 else message,  # Use first message as title
            "created_at": datetime.datetime.now(UTC).isoformat(),
            "updated_at": datetime.datetime.now(UTC).isoformat(),
            "messages": []
        }
    
    # Add user message to conversation
    conversation["messages"].append(user_message)
    conversation["updated_at"] = datetime.datetime.now(UTC).isoformat()
    
    # Save updated conversation
    with open(conversation_file, 'w') as f:
        json.dump(conversation, f)
    
    # Process with Vision API
    try:
        # Read and encode the image file
        with open(image_path, 'rb') as img_file:
            encoded_image = base64.b64encode(img_file.read()).decode('utf-8')
        
        # Get previous conversation messages to add as context
        conversation_history = ""
        if len(conversation["messages"]) > 1:  # If there's more than just the current message
            # Get last 10 messages maximum
            recent_messages = conversation["messages"][-10:] if len(conversation["messages"]) > 10 else conversation["messages"]
            # Format them for context, excluding the most recent (current) message
            conversation_history = "\n".join([
                f"{msg['role'].capitalize()}: {msg['content']}" 
                for msg in recent_messages[:-1]  # Exclude current message
            ])
        
        # Prepare system instruction with history
        system_instruction = bot.get("system_instruction", "You are a helpful assistant that can analyze images.")
        if conversation_history:
            system_instruction += "\n\nThis is the conversation history so far:\n" + conversation_history
            
        # Prepare the message with the image for OpenAI Vision API
        messages = [
            {
                "role": "system", 
                "content": system_instruction
            },
            {
                "role": "user", 
                "content": [
                    {"type": "text", "text": message if message != "[Image uploaded]" else "What's in this image?"},
                    {
                        "type": "image_url",
                        "image_url": {
                            "url": f"data:image/{file_ext};base64,{encoded_image}"
                        }
                    }
                ]
            }
        ]
        
        # Call OpenAI API with vision capabilities
        response = openai.chat.completions.create(
            model="gpt-4o",
            messages=messages,
            max_tokens=1024
        )
        
        response_text = response.choices[0].message.content
        
        # Save the response in the conversation
        bot_response = {
            "id": str(uuid.uuid4()),
            "role": "assistant",
            "content": response_text,
            "timestamp": datetime.datetime.now(UTC).isoformat(),
            "from_image_analysis": True
        }
        conversation["messages"].append(bot_response)
        with open(conversation_file, 'w') as f:
            json.dump(conversation, f)
            
        return jsonify({
            "response": response_text,
            "conversation_id": conversation_id,
            "image_processed": True,
            "debug": {
                "image_path": image_path,
                "message": message,
                "conversation_history_length": len(conversation_history)
            } if debug_mode else None
        }), 200
        
    except Exception as e:
        print(f"Error in chat_with_image: {str(e)}")
        return jsonify({
            "error": "I apologize, but I encountered an error while processing your image. Please try again or contact support if the issue persists.",
            "details": str(e) if debug_mode else None
        }), 500
    
   
# This will serve the React frontend from the bundled location
@app.route('/', defaults={'path': ''})
@app.route('/<path:path>')
def serve_frontend(path):
    # Check if app is running in bundled mode
    if getattr(sys, 'frozen', False) and hasattr(sys, '_MEIPASS'):
        frontend_path = os.path.join(sys._MEIPASS, 'frontend_build')
        if path and os.path.exists(os.path.join(frontend_path, path)):
            return send_from_directory(frontend_path, path)
        return send_from_directory(frontend_path, 'index.html')
    # In development mode, let React dev server handle frontend
    return jsonify({"message": "API endpoint working. Use React dev server for frontend."})

@app.route('/api/datasets/<dataset_id>/documents/<document_id>/download/<filename>', methods=['GET'])
@require_auth
def download_document(user_data, dataset_id, document_id, filename):
    # Check if dataset exists and belongs to user
    datasets_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "datasets")
    user_datasets_file = os.path.join(datasets_dir, f"{user_data['id']}_datasets.json")
    
    if not os.path.exists(user_datasets_file):
        return jsonify({"error": "Dataset not found"}), 404
        
    with open(user_datasets_file, 'r') as f:
        datasets = json.load(f)
        
    dataset_exists = False
    for dataset in datasets:
        if dataset["id"] == dataset_id:
            dataset_exists = True
            break
            
    if not dataset_exists:
        return jsonify({"error": "Dataset not found"}), 404

    # Look for the document in the uploads folder
    document_path = None
    
    # First, try direct match with document_id
    direct_match = os.path.join(DOCUMENT_FOLDER, f"{document_id}_{filename}")
    if os.path.exists(direct_match):
        document_path = direct_match
    else:
        # If direct match fails, search for any file ending with the filename
        for file in os.listdir(DOCUMENT_FOLDER):
            if file.endswith(f"_{filename}"):
                document_path = os.path.join(DOCUMENT_FOLDER, file)
                break

    if not document_path or not os.path.exists(document_path):
        return jsonify({"error": "Document not found"}), 404

    return send_file(document_path, as_attachment=True, download_name=filename)

@app.route('/api/datasets/<dataset_id>/images', methods=['GET'])
@require_auth
def get_dataset_images(user_data, dataset_id):
    """Get all images for a specific dataset"""
    # Check if dataset exists and belongs to user
    datasets_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "datasets")
    user_datasets_file = os.path.join(datasets_dir, f"{user_data['id']}_datasets.json")
    
    if not os.path.exists(user_datasets_file):
        return jsonify({"error": "Dataset not found"}), 404
        
    with open(user_datasets_file, 'r') as f:
        datasets = json.load(f)
        
    dataset_exists = False
    for dataset in datasets:
        if dataset["id"] == dataset_id:
            dataset_exists = True
            break
            
    if not dataset_exists:
        return jsonify({"error": "Dataset not found"}), 404
    
    # Get images metadata from the image processor
    if dataset_id in image_processor.image_metadata:
        images = []
        for img_meta in image_processor.image_metadata[dataset_id]:
            # Create web-accessible URLs for each image
            img_meta_copy = img_meta.copy()
            if 'path' in img_meta_copy:
                img_meta_copy['url'] = f"/api/images/{os.path.basename(img_meta_copy['path'])}"
                # Don't expose the full path in API
                if 'path' in img_meta_copy:
                    del img_meta_copy['path']
            images.append(img_meta_copy)
        
        # Sort images by creation date if available (newest first)
        images.sort(key=lambda x: x.get('created_at', ''), reverse=True)
        
        return jsonify({
            "images": images,
            "total_images": len(images)
        }), 200
    else:
        return jsonify({
            "images": [],
            "total_images": 0
        }), 200

@app.route('/api/datasets/<dataset_id>/images', methods=['POST'])
@require_auth
def upload_image(user_data, dataset_id):
    """Upload an image to a dataset"""
    # Check if dataset exists and belongs to user
    datasets_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "datasets")
    user_datasets_file = os.path.join(datasets_dir, f"{user_data['id']}_datasets.json")
    
    if not os.path.exists(user_datasets_file):
        return jsonify({"error": "Dataset not found"}), 404
        
    with open(user_datasets_file, 'r') as f:
        datasets = json.load(f)
        
    dataset_exists = False
    dataset_index = -1
    for idx, dataset in enumerate(datasets):
        if dataset["id"] == dataset_id:
            dataset_exists = True
            dataset_index = idx
            break
            
    if not dataset_exists:
        return jsonify({"error": "Dataset not found"}), 404
    
    # Handle both form data and JSON requests
    if request.content_type and request.content_type.startswith('multipart/form-data'):
        # Handle image file from form data
        if 'image' not in request.files:
            return jsonify({"error": "No image file provided"}), 400
            
        image_file = request.files['image']
        if image_file.filename == '':
            return jsonify({"error": "No image selected"}), 400
            
        # Save the image
        filename = secure_filename(image_file.filename)
        image_path = os.path.join(IMAGE_FOLDER, f"{uuid.uuid4()}_{filename}")
        image_file.save(image_path)
        
    else:
        # Handle base64 image data from JSON
        data = request.json
        if not data or 'image_data' not in data:
            return jsonify({"error": "Image data is required"}), 400
            
        image_data = data['image_data']
        filename = data.get('filename', f"image_{uuid.uuid4()}.jpg")
        
        # Extract the base64 content
        try:
            if ';base64,' in image_data:
                header, encoded = image_data.split(';base64,')
                file_ext = header.split('/')[-1]
                if not file_ext:
                    file_ext = 'jpg'
            else:
                encoded = image_data
                file_ext = 'jpg'
                
            # Generate a unique filename
            if not filename:
                filename = f"image_{uuid.uuid4()}.{file_ext}"
                
            # Save the image
            image_path = os.path.join(IMAGE_FOLDER, secure_filename(filename))
            with open(image_path, 'wb') as f:
                f.write(base64.b64decode(encoded))
                
        except Exception as e:
            return jsonify({"error": f"Failed to process image data: {str(e)}"}), 400
    
    # Resize image if it's too large
    try:
        image_path = resize_image(image_path)
    except Exception as e:
        print(f"Warning: Failed to resize image: {str(e)}")
    
    # Get additional metadata
    custom_metadata = {}
    if request.content_type and request.content_type.startswith('multipart/form-data'):
        # Get metadata from form fields
        custom_metadata['description'] = request.form.get('description', '')
        custom_metadata['tags'] = request.form.get('tags', '').split(',') if request.form.get('tags') else []
    else:
        # Get metadata from JSON
        custom_metadata['description'] = data.get('description', '')
        custom_metadata['tags'] = data.get('tags', [])
    
    # Add user information to metadata
    custom_metadata['user_id'] = user_data['id']
    custom_metadata['username'] = user_data['username']
    
    try:
        # Add image to dataset in image processor
        image_metadata = image_processor.add_image_to_dataset(dataset_id, image_path, custom_metadata)
        
        # Update image count in dataset
        datasets[dataset_index]["image_count"] = datasets[dataset_index].get("image_count", 0) + 1
        with open(user_datasets_file, 'w') as f:
            json.dump(datasets, f)
            
        # Generate URL for the image
        image_url = f"/api/images/{os.path.basename(image_path)}"
        
        return jsonify({
            "message": "Image uploaded and processed successfully",
            "image": {
                "id": image_metadata["id"],
                "filename": os.path.basename(image_path),
                "url": image_url,
                "caption": image_metadata["caption"],
                "description": custom_metadata.get("description", ""),
                "tags": custom_metadata.get("tags", [])
            }
        }), 200
        
    except Exception as e:
        print(f"Error processing image: {str(e)}")
        if os.path.exists(image_path):
            os.remove(image_path)
        return jsonify({"error": f"Failed to process image: {str(e)}"}), 500

@app.route('/api/datasets/<dataset_id>/images/<image_id>', methods=['DELETE'])
@require_auth
def remove_image(user_data, dataset_id, image_id):
    """Remove an image from a dataset"""
    # Check if dataset exists and belongs to user
    datasets_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "datasets")
    user_datasets_file = os.path.join(datasets_dir, f"{user_data['id']}_datasets.json")
    
    if not os.path.exists(user_datasets_file):
        return jsonify({"error": "Dataset not found"}), 404
        
    with open(user_datasets_file, 'r') as f:
        datasets = json.load(f)
        
    dataset_exists = False
    dataset_index = -1
    for idx, dataset in enumerate(datasets):
        if dataset["id"] == dataset_id:
            dataset_exists = True
            dataset_index = idx
            break
            
    if not dataset_exists:
        return jsonify({"error": "Dataset not found"}), 404
    
    try:
        # Get image path before removal for cleanup
        image_path = None
        if dataset_id in image_processor.image_metadata:
            for img in image_processor.image_metadata[dataset_id]:
                if img["id"] == image_id:
                    image_path = img["path"]
                    break
        
        # Remove image from dataset
        success = image_processor.remove_image(dataset_id, image_id)
        if not success:
            return jsonify({"error": "Image not found in dataset"}), 404
            
        # Delete the image file if it exists
        if image_path and os.path.exists(image_path):
            try:
                os.remove(image_path)
            except Exception as e:
                print(f"Warning: Failed to delete image file: {str(e)}")
        
        # Update image count in dataset
        if "image_count" in datasets[dataset_index] and datasets[dataset_index]["image_count"] > 0:
            datasets[dataset_index]["image_count"] -= 1
        else:
            # Set image count to actual count if it wasn't tracked before
            image_count = 0
            if dataset_id in image_processor.image_metadata:
                image_count = len(image_processor.image_metadata[dataset_id])
            datasets[dataset_index]["image_count"] = image_count
            
        with open(user_datasets_file, 'w') as f:
            json.dump(datasets, f)
                
        return jsonify({"message": "Image removed successfully"}), 200
        
    except Exception as e:
        print(f"Error removing image: {str(e)}")
        return jsonify({"error": f"Failed to remove image: {str(e)}"}), 500

@app.route('/api/datasets/<dataset_id>/search-images', methods=['POST'])
@require_auth
def search_dataset_images(user_data, dataset_id):
    """Search for images in a dataset using a text query"""
    # Check if dataset exists and belongs to user
    datasets_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "datasets")
    user_datasets_file = os.path.join(datasets_dir, f"{user_data['id']}_datasets.json")
    
    if not os.path.exists(user_datasets_file):
        return jsonify({"error": "Dataset not found"}), 404
        
    with open(user_datasets_file, 'r') as f:
        datasets = json.load(f)
        
    dataset_exists = False
    for dataset in datasets:
        if dataset["id"] == dataset_id:
            dataset_exists = True
            break
            
    if not dataset_exists:
        return jsonify({"error": "Dataset not found"}), 404
    
    # Get query parameters
    data = request.json
    if not data or not data.get('query'):
        return jsonify({"error": "Search query is required"}), 400
        
    query = data['query']
    limit = data.get('limit', 5)  # Default to 5 results
    
    try:
        # Search for images
        results = image_processor.search_images(dataset_id, query, limit)
        
        # Format results
        formatted_results = []
        for result in results:
            result_copy = result.copy()
            if 'path' in result_copy:
                result_copy['url'] = f"/api/images/{os.path.basename(result_copy['path'])}"
                del result_copy['path']  # Don't expose full path
            formatted_results.append(result_copy)
            
        return jsonify({
            "query": query,
            "results": formatted_results
        }), 200
        
    except Exception as e:
        print(f"Error searching images: {str(e)}")
        return jsonify({"error": f"Failed to search images: {str(e)}"}), 500

# Helper function to resize large images
def resize_image(image_path, max_size=1024):
    """Resize an image if it's larger than max_size in either dimension
    
    Args:
        image_path: Path to the image file
        max_size: Maximum size for either dimension
        
    Returns:
        str: Path to the resized image (same as input if no resize needed)
    """
    try:
        with Image.open(image_path) as img:
            width, height = img.size
            
            # Check if resize needed
            if width <= max_size and height <= max_size:
                return image_path
                
            # Calculate new dimensions maintaining aspect ratio
            if width > height:
                new_width = max_size
                new_height = int(height * (max_size / width))
            else:
                new_height = max_size
                new_width = int(width * (max_size / height))
                
            # Resize the image
            resized = img.resize((new_width, new_height), Image.Resampling.LANCZOS)
            
            # Save the resized image (overwrite original)
            resized.save(image_path, optimize=True, quality=85)
            print(f"Resized image from {width}x{height} to {new_width}x{new_height}")
            
            return image_path
    except Exception as e:
        print(f"Error resizing image: {str(e)}")
        return image_path  # Return original path if resize fails

@app.route('/api/datasets/<dataset_id>/type', methods=['GET'])
@require_auth
def get_dataset_type(user_data, dataset_id):
    """Get the type of a dataset (text, image, or mixed) to inform frontend file selection"""
    # Check if dataset exists and belongs to user
    datasets_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "datasets")
    user_datasets_file = os.path.join(datasets_dir, f"{user_data['id']}_datasets.json")
    
    if not os.path.exists(user_datasets_file):
        return jsonify({"error": "Dataset not found"}), 404
        
    with open(user_datasets_file, 'r') as f:
        datasets = json.load(f)
        
    for dataset in datasets:
        if dataset["id"] == dataset_id:
            # Return the dataset type and accepted file types
            dataset_type = dataset.get("type", "text")
            accepted_files = {
                "text": [".pdf", ".docx", ".txt", ".pptx"],
                "image": [".jpg", ".jpeg", ".png", ".gif", ".webp", ".bmp"],
                "mixed": [".pdf", ".docx", ".txt", ".pptx", ".jpg", ".jpeg", ".png", ".gif", ".webp", ".bmp"]
            }
            
            return jsonify({
                "type": dataset_type,
                "supported_extensions": accepted_files.get(dataset_type, []),
                "mime_types": get_mime_types_for_dataset(dataset_type)
            }), 200
            
    return jsonify({"error": "Dataset not found"}), 404

def get_mime_types_for_dataset(dataset_type):
    """Helper function to return appropriate MIME types for dataset type"""
    mime_types = {
        "text": "application/pdf,application/vnd.openxmlformats-officedocument.wordprocessingml.document,text/plain,application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "image": "image/jpeg,image/png,image/gif,image/webp,image/bmp",
        "mixed": "application/pdf,application/vnd.openxmlformats-officedocument.wordprocessingml.document,text/plain,application/vnd.openxmlformats-officedocument.presentationml.presentation,image/jpeg,image/png,image/gif,image/webp,image/bmp"
    }
    
    return mime_types.get(dataset_type, "")

@app.route('/api/upload-example', methods=['GET'])
def upload_example():
    """Serve a simple HTML example for uploading files that supports both text and image files"""
    html = """
    <!DOCTYPE html>
    <html>
    <head>
        <title>RagBot File Upload Example</title>
        <style>
            body { font-family: Arial, sans-serif; max-width: 800px; margin: 0 auto; padding: 20px; }
            .container { border: 1px solid #ccc; border-radius: 5px; padding: 20px; margin-bottom: 20px; }
            label { display: block; margin-bottom: 10px; font-weight: bold; }
            select, input, button { margin-bottom: 10px; padding: 8px; width: 100%; }
            button { background-color: #4CAF50; color: white; border: none; cursor: pointer; }
            button:hover { background-color: #45a049; }
            .result { margin-top: 20px; padding: 10px; background-color: #f5f5f5; border-radius: 5px; white-space: pre-wrap; }
            img { max-width: 100%; max-height: 300px; display: block; margin: 10px 0; }
        </style>
    </head>
    <body>
        <h1>RagBot File Upload Example</h1>
        
        <div class="container">
            <h2>Step 1: Select Dataset</h2>
            <label for="dataset">Choose a dataset:</label>
            <select id="dataset">
                <option value="">Loading datasets...</option>
            </select>
            <div id="dataset-info"></div>
        </div>
        
        <div class="container">
            <h2>Step 2: Upload File</h2>
            <div id="upload-form" style="display: none;">
                <label for="file">Select a file:</label>
                <input type="file" id="file" name="file">
                
                <label for="description">Description (optional):</label>
                <input type="text" id="description" name="description" placeholder="Enter a description">
                
                <label for="tags">Tags (comma-separated, optional):</label>
                <input type="text" id="tags" name="tags" placeholder="tag1, tag2, tag3">
                
                <button id="upload">Upload</button>
            </div>
        </div>
        
        <div class="container">
            <h2>Result</h2>
            <div id="result" class="result">No upload yet</div>
            <div id="image-preview"></div>
        </div>
        
        <script>
            // Get the authentication token from localStorage or prompt the user
            let token = localStorage.getItem('authToken');
            if (!token) {
                token = prompt('Please enter your authentication token:');
                if (token) localStorage.setItem('authToken', token);
            }
            
            // Fetch all datasets
            async function fetchDatasets() {
                try {
                    const response = await fetch('/api/datasets', {
                        headers: {
                            'Authorization': `Bearer ${token}`
                        }
                    });
                    
                    if (!response.ok) throw new Error('Failed to fetch datasets');
                    
                    const datasets = await response.json();
                    const select = document.getElementById('dataset');
                    select.innerHTML = '';
                    
                    if (datasets.length === 0) {
                        select.innerHTML = '<option value="">No datasets available</option>';
                        return;
                    }
                    
                    datasets.forEach(dataset => {
                        const option = document.createElement('option');
                        option.value = dataset.id;
                        option.textContent = `${dataset.name} (${dataset.type || 'text'})`;
                        option.dataset.type = dataset.type || 'text';
                        select.appendChild(option);
                    });
                    
                    // Show dataset info and setup file input
                    updateDatasetInfo();
                } catch (error) {
                    document.getElementById('result').textContent = `Error: ${error.message}`;
                }
            }
            
            // Update dataset info and file input accept attribute
            async function updateDatasetInfo() {
                const select = document.getElementById('dataset');
                const datasetId = select.value;
                
                if (!datasetId) {
                    document.getElementById('dataset-info').textContent = 'Please select a dataset';
                    document.getElementById('upload-form').style.display = 'none';
                    return;
                }
                
                try {
                    const response = await fetch(`/api/datasets/${datasetId}/type`, {
                        headers: {
                            'Authorization': `Bearer ${token}`
                        }
                    });
                    
                    if (!response.ok) throw new Error('Failed to fetch dataset info');
                    
                    const info = await response.json();
                    const datasetInfo = document.getElementById('dataset-info');
                    const fileInput = document.getElementById('file');
                    
                    datasetInfo.textContent = `Type: ${info.type}. Supported files: ${info.supported_extensions.join(', ')}`;
                    fileInput.accept = info.mime_types;
                    
                    document.getElementById('upload-form').style.display = 'block';
                } catch (error) {
                    document.getElementById('dataset-info').textContent = `Error: ${error.message}`;
                }
            }
            
            // Handle file upload
            async function uploadFile() {
                const datasetId = document.getElementById('dataset').value;
                const fileInput = document.getElementById('file');
                const description = document.getElementById('description').value;
                const tags = document.getElementById('tags').value;
                const resultDiv = document.getElementById('result');
                const imagePreview = document.getElementById('image-preview');
                
                if (!datasetId) {
                    resultDiv.textContent = 'Please select a dataset';
                    return;
                }
                
                if (!fileInput.files || fileInput.files.length === 0) {
                    resultDiv.textContent = 'Please select a file to upload';
                    return;
                }
                
                resultDiv.textContent = 'Uploading...';
                imagePreview.innerHTML = '';
                
                const formData = new FormData();
                formData.append('file', fileInput.files[0]);
                if (description) formData.append('description', description);
                if (tags) formData.append('tags', tags);
                
                try {
                    const response = await fetch(`/api/datasets/${datasetId}/documents`, {
                        method: 'POST',
                        headers: {
                            'Authorization': `Bearer ${token}`
                        },
                        body: formData
                    });
                    
                    const result = await response.json();
                    
                    if (!response.ok) {
                        resultDiv.textContent = `Error: ${result.error}`;
                        return;
                    }
                    
                    resultDiv.textContent = JSON.stringify(result, null, 2);
                    
                    // If the response contains an image, display it
                    if (result.image && result.image.url) {
                        const img = document.createElement('img');
                        img.src = result.image.url;
                        img.alt = result.image.caption || 'Uploaded image';
                        imagePreview.appendChild(img);
                        
                        const caption = document.createElement('p');
                        caption.textContent = `Caption: ${result.image.caption || 'No caption available'}`;
                        imagePreview.appendChild(caption);
                    }
                } catch (error) {
                    resultDiv.textContent = `Error: ${error.message}`;
                }
            }
            
            // Setup event listeners
            document.addEventListener('DOMContentLoaded', () => {
                fetchDatasets();
                
                document.getElementById('dataset').addEventListener('change', updateDatasetInfo);
                document.getElementById('upload').addEventListener('click', uploadFile);
            });
        </script>
    </body>
    </html>
    """
    return html

@app.route('/api/documents/<document_id>/content', methods=['GET'])
@require_auth
def get_document_content(user_data, document_id):
    """Get the full content of a document"""
    try:
        # Try to find the document in ChromaDB
        # First we need to search across all datasets the user has access to
        datasets_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "datasets")
        user_datasets_file = os.path.join(datasets_dir, f"{user_data['id']}_datasets.json")
        
        if not os.path.exists(user_datasets_file):
            return jsonify({"error": "No datasets found"}), 404
            
        with open(user_datasets_file, 'r') as f:
            datasets = json.load(f)
        
        # Search for the document in each dataset
        document_chunks = []
        document_metadata = None
        
        for dataset in datasets:
            dataset_id = dataset["id"]
            try:
                collection = chroma_client.get_collection(name=dataset_id, embedding_function=openai_ef)
                results = collection.get(where={"document_id": document_id})
                
                if results and len(results["documents"]) > 0:
                    # Found chunks from this document
                    document_chunks.extend(list(zip(results["documents"], results["metadatas"])))
                    # Get metadata from the first chunk
                    if not document_metadata and results["metadatas"] and len(results["metadatas"]) > 0:
                        document_metadata = results["metadatas"][0]
            except Exception as e:
                print(f"Error querying collection {dataset_id}: {str(e)}")
                continue
        
        if not document_chunks:
            return jsonify({"error": "Document not found"}), 404
            
        # Sort chunks by their position in the document
        document_chunks.sort(key=lambda x: x[1].get("chunk", 0) if x[1] else 0)
        
        # Combine all chunks into a single document
        full_content = "\n\n".join([chunk[0] for chunk in document_chunks])
        
        # Get document metadata
        filename = document_metadata.get("filename", "Unknown document") if document_metadata else "Unknown document"
        file_path = document_metadata.get("file_path", "") if document_metadata else ""
        
        # Try to get the original file if it exists
        original_content = ""
        if file_path and os.path.exists(file_path):
            try:
                original_content = extract_text_from_file(file_path)
            except Exception as e:
                print(f"Error extracting text from original file: {str(e)}")
                # Use the reconstructed content from chunks as fallback
                original_content = full_content
        else:
            # Use the reconstructed content from chunks as fallback
            original_content = full_content
        
        return jsonify({
            "document_id": document_id,
            "filename": filename,
            "content": original_content if original_content else full_content
        }), 200
        
    except Exception as e:
        print(f"Error retrieving document content: {str(e)}")
        return jsonify({"error": f"Failed to retrieve document content: {str(e)}"}), 500

@app.route('/api/context/<document_id>', methods=['GET'])
@require_auth
def get_context_snippet(user_data, document_id):
    """Get the full content of a context snippet with option to view the entire document"""
    try:
        # Get snippet index from query parameter, if provided
        chunk_index = request.args.get('chunk', None)
        
        # Try to find the document in ChromaDB
        # First we need to search across all datasets the user has access to
        datasets_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "datasets")
        user_datasets_file = os.path.join(datasets_dir, f"{user_data['id']}_datasets.json")
        
        if not os.path.exists(user_datasets_file):
            return jsonify({"error": "No datasets found"}), 404
            
        with open(user_datasets_file, 'r') as f:
            datasets = json.load(f)
        
        # Search for the document in each dataset
        document_chunks = []
        document_metadata = None
        
        for dataset in datasets:
            dataset_id = dataset["id"]
            try:
                collection = chroma_client.get_collection(name=dataset_id, embedding_function=openai_ef)
                
                # If chunk index is provided, get only that specific chunk
                if chunk_index:
                    chunk_id = f"{document_id}_{chunk_index}"
                    results = collection.get(ids=[chunk_id])
                else:
                    # Otherwise get all chunks for this document
                    results = collection.get(where={"document_id": document_id})
                
                if results and len(results["documents"]) > 0:
                    # Found chunks from this document
                    document_chunks.extend(list(zip(results["documents"], results["metadatas"])))
                    # Get metadata from the first chunk
                    if not document_metadata and results["metadatas"] and len(results["metadatas"]) > 0:
                        document_metadata = results["metadatas"][0]
            except Exception as e:
                print(f"Error querying collection {dataset_id}: {str(e)}")
                continue
        
        if not document_chunks:
            return jsonify({"error": "Document chunk not found"}), 404
            
        # Sort chunks by their position in the document
        document_chunks.sort(key=lambda x: x[1].get("chunk", 0) if x[1] else 0)
        
        # Get snippet content - if specific chunk requested, return just that one
        if chunk_index:
            snippet_content = document_chunks[0][0] if document_chunks else ""
        else:
            # Otherwise combine all chunks
            snippet_content = "\n\n".join([chunk[0] for chunk in document_chunks])
        
        # Get document metadata
        filename = document_metadata.get("filename", "Unknown document") if document_metadata else "Unknown document"
        file_path = document_metadata.get("file_path", "") if document_metadata else ""
        source = document_metadata.get("source", filename) if document_metadata else filename
        
        # Check if the original file exists - will be used by frontend to offer "view original" option
        original_file_exists = file_path and os.path.exists(file_path)
        
        return jsonify({
            "document_id": document_id,
            "filename": filename,
            "source": source,
            "content": snippet_content,
            "original_file_exists": original_file_exists,
            "chunk_index": chunk_index
        }), 200
        
    except Exception as e:
        print(f"Error retrieving context snippet: {str(e)}")
        return jsonify({"error": f"Failed to retrieve context snippet: {str(e)}"}), 500

@app.route('/api/documents/<document_id>/original', methods=['GET'])
@require_auth
def get_original_document(user_data, document_id):
    """Get the original file content if it exists"""
    try:
        # Try to find the document metadata in ChromaDB
        datasets_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "datasets")
        user_datasets_file = os.path.join(datasets_dir, f"{user_data['id']}_datasets.json")
        
        if not os.path.exists(user_datasets_file):
            return jsonify({"error": "No datasets found"}), 404
            
        with open(user_datasets_file, 'r') as f:
            datasets = json.load(f)
        
        # Search for the document in each dataset to get its metadata
        document_metadata = None
        
        for dataset in datasets:
            dataset_id = dataset["id"]
            try:
                collection = chroma_client.get_collection(name=dataset_id, embedding_function=openai_ef)
                
                # Query for just one chunk to get metadata
                results = collection.get(
                    where={"document_id": document_id},
                    limit=1
                )
                
                if results and len(results["metadatas"]) > 0:
                    document_metadata = results["metadatas"][0]
                    break
            except Exception as e:
                print(f"Error querying collection {dataset_id}: {str(e)}")
                continue
        
        if not document_metadata:
            return jsonify({"error": "Document not found"}), 404
            
        # Get file path from metadata
        file_path = document_metadata.get("file_path", "")
        filename = document_metadata.get("filename", "Unknown document")
        
        # Check if file exists
        if not file_path or not os.path.exists(file_path):
            return jsonify({"error": "Original file not found"}), 404
            
        # Extract text from the original file
        try:
            content = extract_text_from_file(file_path)
            file_extension = os.path.splitext(file_path)[1].lower()
            
            return jsonify({
                "document_id": document_id,
                "filename": filename,
                "content": content,
                "file_type": file_extension[1:] if file_extension.startswith('.') else file_extension
            }), 200
        except Exception as e:
            return jsonify({"error": f"Error extracting text from file: {str(e)}"}), 500
            
    except Exception as e:
        print(f"Error retrieving original document: {str(e)}")
        return jsonify({"error": f"Failed to retrieve original document: {str(e)}"}), 500

if __name__ == '__main__':
    app.run(host='0.0.0.0', port=5000, debug=True)
