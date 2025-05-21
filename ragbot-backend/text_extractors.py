import os
import base64
import io
import uuid
import re
import numpy as np
from PIL import Image
import pytesseract
import fitz  # PyMuPDF
from pdf2image import convert_from_path
from PyPDF2 import PdfReader
import docx
from pptx import Presentation

# Configure OCR
try:
    # On Windows, you might need to set the path to Tesseract executable
    if os.name == 'nt':  # If on Windows
        pytesseract.pytesseract.tesseract_cmd = os.getenv('TESSERACT_PATH', r'C:\Program Files\Tesseract-OCR\tesseract.exe')
except Exception as e:
    print(f"Warning: Could not configure pytesseract: {str(e)}")

# Define a function to extract text from images using OCR
def extract_text_from_image(img):
    """Extract text from an image using OCR
    
    Args:
        img: PIL Image object or path to image
        
    Returns:
        str: Extracted text
    """
    try:
        if isinstance(img, str):
            # If img is a path
            img = Image.open(img)
        
        # Convert image to RGB if it's not
        if img.mode != 'RGB':
            img = img.convert('RGB')
            
        # Apply image preprocessing to improve OCR
        # Convert to grayscale
        gray = img.convert('L')
        
        # Apply threshold to create a black and white image
        threshold = 128
        bw = gray.point(lambda x: 0 if x < threshold else 255, '1')
        
        # Extract text using pytesseract
        text = pytesseract.image_to_string(bw)
        
        return text
    except Exception as e:
        print(f"Warning: Failed to extract text from image: {str(e)}")
        return ""

def extract_text_from_pdf(pdf_path):
    """Extract text from a PDF including text in images
    
    Args:
        pdf_path: Path to the PDF file
        
    Returns:
        str: Extracted text
    """
    text = ""
    image_text = ""
    
    # First, extract regular text using PyPDF2
    with open(pdf_path, 'rb') as file:
        reader = PdfReader(file)
        for page in reader.pages:
            text += page.extract_text() + "\n"
    
    # Then, extract text from images using PyMuPDF and OCR
    try:
        doc = fitz.open(pdf_path)
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            
            # Get images from the page
            image_list = page.get_images(full=True)
            
            # Process each image
            for img_index, img in enumerate(image_list):
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                
                # Load image using PIL
                img = Image.open(io.BytesIO(image_bytes))
                
                # Extract text using OCR
                img_text = extract_text_from_image(img)
                if img_text.strip():
                    image_text += f"Image {page_num+1}.{img_index+1} text: {img_text}\n"
        
        doc.close()
    except Exception as e:
        print(f"Warning: Failed to extract images from PDF: {str(e)}")
        
        # Fall back to pdf2image if PyMuPDF fails
        try:
            # Convert PDF to images
            images = convert_from_path(pdf_path)
            
            # Process each page image
            for i, img in enumerate(images):
                # Extract text using OCR
                img_text = extract_text_from_image(img)
                if img_text.strip():
                    image_text += f"Page {i+1} image text: {img_text}\n"
        except Exception as e2:
            print(f"Warning: Failed to use pdf2image fallback: {str(e2)}")
    
    # Combine regular text and image text
    combined_text = text
    if image_text:
        combined_text += "\n\nExtracted from images:\n" + image_text
        
    return combined_text

def extract_text_from_docx(docx_path):
    """Extract text from a Word document
    
    Args:
        docx_path: Path to the docx file
        
    Returns:
        str: Extracted text
    """
    text = ""
    doc = docx.Document(docx_path)
    
    # Extract text from paragraphs
    for para in doc.paragraphs:
        text += para.text + "\n"
        
    # Extract text from tables
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                text += cell.text + " "
            text += "\n"
    
    # Unfortunately, direct image extraction from docx is complex
    # We don't have a straightforward way to extract images for OCR
    # A complete solution would require unpacking the docx and processing embedded images
    
    return text

def extract_text_from_pptx(pptx_path):
    """Extract text from a PowerPoint presentation including text in images
    
    Args:
        pptx_path: Path to the pptx file
        
    Returns:
        tuple: (Extracted text, List of image metadata dicts)
    """
    summary_lines = []
    image_text = ""
    image_captions = ""
    prs = Presentation(pptx_path)
    slide_image_metadata = []
    
    # Log the total number of slides
    slide_count = len(prs.slides)
    print(f"Processing PowerPoint with {slide_count} slides")
    
    # Add presentation metadata to summary
    summary_lines.append("## PRESENTATION SUMMARY ##")
    summary_lines.append(f"Total Slides: {slide_count}")
    summary_lines.append("")
    
    # Get or create ImageProcessor for captioning
    try:
        from image_processor import ImageProcessor
        img_processor = ImageProcessor(os.path.join(os.path.dirname(os.path.abspath(__file__)), "data"))
    except Exception as e:
        print(f"Warning: Could not initialize ImageProcessor: {str(e)}")
        img_processor = None
    
    # Create temp directory for extracting images
    import tempfile
    with tempfile.TemporaryDirectory() as temp_dir:
        # Process each slide
        for slide_num, slide in enumerate(prs.slides):
            # Add clear slide separator with prominent marker
            summary_lines.append(f"## SLIDE {slide_num+1} ##")
            slide_has_content = False
            
            # Try to get slide title
            slide_title = ""
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text and hasattr(shape, 'is_title') and shape.is_title:
                    slide_title = shape.text_frame.text.strip()
                    summary_lines.append(f"TITLE: {slide_title}")
                    slide_has_content = True
                    break
            
            # Add a section for text content
            text_content = []
            
            # Get text from all shapes with text
            for shape in slide.shapes:
                # Get text from text frames
                if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    if shape.text_frame and shape.text_frame.text and shape.text_frame.text.strip():
                        # Skip if we already included this as the title
                        text_content_item = shape.text_frame.text.strip()
                        if text_content_item != slide_title:
                            text_content.append(text_content_item)
                            slide_has_content = True
                # Also check for text attribute (some shapes have this instead)
                elif hasattr(shape, "text") and shape.text and shape.text.strip():
                    # Skip if we already included this as the title
                    text_content_item = shape.text.strip()
                    if text_content_item != slide_title:
                        text_content.append(text_content_item)
                        slide_has_content = True
                
                # Get text from tables if present
                if shape.has_table:
                    table_content = ["TABLE CONTENT:"]
                    has_table_content = False
                    for row in shape.table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text_frame.text:
                                row_text.append(cell.text_frame.text.strip())
                                has_table_content = True
                        if row_text:
                            table_content.append(" | ".join(row_text))
                    
                    if has_table_content:
                        text_content.extend(table_content)
                        slide_has_content = True
            
            # Add text content section if there's any content
            if text_content:
                summary_lines.append("\nCONTENT:")
                for item in text_content:
                    summary_lines.append(item)
            
            # Add a section for image content
            slide_images = []
            
            # Extract image if this shape is a picture
            for shape_index, shape in enumerate(slide.shapes):
                if hasattr(shape, "shape_type") and shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    try:
                        image_path = os.path.join(temp_dir, f"slide_{slide_num+1}_img_{uuid.uuid4()}.png")
                        with open(image_path, 'wb') as f:
                            f.write(shape.image.blob)
                        image_info = {
                            "slide_number": slide_num+1,
                            "slide_title": slide_title,
                            "filename": os.path.basename(pptx_path),
                            "image_path": image_path,
                            "ocr_text": "",
                            "caption": ""
                        }
                        img_text = extract_text_from_image(image_path)
                        if img_text.strip():
                            image_info["ocr_text"] = img_text.strip()
                            slide_has_content = True
                        if img_processor:
                            try:
                                caption = img_processor.generate_caption(image_path)
                                if caption and caption != "No caption available":
                                    image_info["caption"] = caption
                                    slide_has_content = True
                            except Exception as e:
                                print(f"Warning: Failed to generate caption for image on slide {slide_num+1}: {str(e)}")
                        slide_images.append(image_info)
                        slide_image_metadata.append(image_info)
                    except Exception as e:
                        print(f"Warning: Failed to extract image from slide {slide_num+1}: {str(e)}")
            
            # Add image content to summary
            if slide_images:
                summary_lines.append("\nIMAGES:")
                for i, img_info in enumerate(slide_images):
                    summary_lines.append(f"Image {i+1}:")
                    if img_info["caption"]:
                        summary_lines.append(f"  Description: {img_info['caption']}")
                    if img_info["ocr_text"]:
                        summary_lines.append(f"  Text content: {img_info['ocr_text']}")
            
            # If slide has no content, mark it
            if not slide_has_content:
                summary_lines.append("No content on this slide")
            
            # Add a separator between slides (double newline)
            summary_lines.append("\n")
    
    # Join all the summarized content
    combined_text = "\n".join(summary_lines)
    
    # Log the overall extraction result
    print(f"PowerPoint extraction complete: {len(combined_text)} characters extracted")
    if not combined_text.strip():
        print("WARNING: No content was extracted from the PowerPoint file!")
        
    return combined_text, slide_image_metadata

def extract_text_from_file(file_path):
    """Extract text from a file based on its extension
    
    Args:
        file_path: Path to the file
        
    Returns:
        str: Extracted text
    """
    extension = os.path.splitext(file_path)[1].lower()
    if extension == '.pdf':
        return extract_text_from_pdf(file_path)
    elif extension == '.docx':
        return extract_text_from_docx(file_path)
    elif extension == '.pptx':
        return extract_text_from_pptx(file_path)
    elif extension == '.txt':
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    else:
        return ""

def create_semantic_chunks(text, max_chunk_size=1000, overlap=200):
    """Create semantically coherent chunks from text, respecting document structure
    
    Args:
        text: The text to chunk
        max_chunk_size: Maximum size of a chunk in characters
        overlap: Number of characters to overlap between chunks
        
    Returns:
        list: List of text chunks
    """
    # Detect if content is from a PowerPoint presentation
    is_powerpoint = "## SLIDE " in text or "## PRESENTATION SUMMARY ##" in text
    
    # Use larger chunks for PowerPoint content
    if is_powerpoint:
        max_chunk_size = 3000  # Much larger to contain multiple slides
        overlap = 500  # More overlap to maintain context between chunks
    
    # If text is empty or too short, return as is
    if not text or len(text) <= max_chunk_size:
        return [text]
    
    # For PowerPoint content, use a special chunking strategy
    if is_powerpoint:
        return chunk_powerpoint_content(text, max_chunk_size, overlap)
    
    # Split the text into sections based on headers or meaningful separators
    # Look for common headers like "Chapter", "Section", or numbered headers
    section_patterns = [
        r'(?:\n|\r\n|\r)#{1,6}\s+(.+?)(?:\n|\r\n|\r)',  # Markdown headers
        r'(?:\n|\r\n|\r)([A-Z][^.\n]{0,50}:)(?:\n|\r\n|\r)',  # Capitalized section names with colon
        r'(?:\n|\r\n|\r)((?:Section|Chapter|Part)\s+\d+[.:]\s*[^\n]+)(?:\n|\r\n|\r)',  # Section/Chapter headings
        r'(?:\n|\r\n|\r)(\d+\.\d*\s+[^\n]{0,50})(?:\n|\r\n|\r)',  # Numbered sections (e.g., "1.2 Title")
        r'(?:\n|\r\n|\r)(Slide \d+:)',  # Slide markers
        r'(?:\n|\r\n|\r)(Title: .*?)(?:\n|\r\n|\r)',  # PowerPoint slide titles
        r'(?:\n|\r\n|\r)(Speaker Notes: .*?)(?:\n|\r\n|\r)',  # PowerPoint speaker notes
        r'(?:\n|\r\n|\r)(Image \d+\.\d+ text:)',  # Image text markers
        r'(?:\n|\r\n|\r)(Page \d+ image text:)',  # Page image text markers
        r'(?:\n|\r\n|\r)(Extracted from images:)',  # OCR section marker
        r'(?:\n|\r\n|\r)(Presentation Title: .*?)(?:\n|\r\n|\r)',  # Presentation title
        r'(?:\n|\r\n|\r)(Author: .*?)(?:\n|\r\n|\r)'  # Author info
    ]
    
    # Find all potential section boundaries
    sections = []
    last_end = 0
    
    # First, identify all possible section boundaries
    boundaries = []
    for pattern in section_patterns:
        for match in re.finditer(pattern, text, re.MULTILINE):
            boundaries.append((match.start(), match.group(0)))
    
    # Sort boundaries by position
    boundaries.sort(key=lambda x: x[0])
    
    # Create sections based on boundaries
    for i, (pos, header) in enumerate(boundaries):
        if pos > last_end:  # Only process if we haven't already included this position
            if i < len(boundaries) - 1:
                next_pos = boundaries[i+1][0]
                section_text = text[pos:next_pos]
                sections.append(section_text)
                last_end = next_pos
            else:
                # Last section
                section_text = text[pos:]
                sections.append(section_text)
                last_end = len(text)
    
    # If no sections were found or first section doesn't start at beginning
    if not sections or boundaries[0][0] > 0 if boundaries else True:
        if boundaries:
            first_pos = boundaries[0][0]
            sections.insert(0, text[:first_pos])
        else:
            # If no sections were identified, treat the whole text as one section
            sections = [text]
    
    # Further chunk each section if it's still too large
    chunks = []
    for section in sections:
        if len(section) <= max_chunk_size:
            chunks.append(section)
        else:
            # Try to split on paragraph boundaries
            paragraphs = re.split(r'(?:\n|\r\n|\r){2,}', section)
            current_chunk = ""
            
            for para in paragraphs:
                # If adding this paragraph exceeds max size and we already have content
                if len(current_chunk) + len(para) > max_chunk_size and current_chunk:
                    chunks.append(current_chunk)
                    # Start new chunk with overlap
                    overlap_point = max(0, len(current_chunk) - overlap)
                    current_chunk = current_chunk[overlap_point:] + para
                else:
                    current_chunk += para
                    
                if len(current_chunk) > max_chunk_size:
                    # If a single paragraph is too big, we have to split within the paragraph
                    sentences = re.split(r'(?<=[.!?])\s+', current_chunk)
                    current_chunk = ""
                    
                    for sentence in sentences:
                        if len(current_chunk) + len(sentence) > max_chunk_size and current_chunk:
                            chunks.append(current_chunk)
                            # Start new chunk with overlap
                            overlap_point = max(0, len(current_chunk) - overlap)
                            current_chunk = current_chunk[overlap_point:] + sentence
                        else:
                            current_chunk += sentence + " "
            
            # Don't forget to add the last chunk
            if current_chunk:
                chunks.append(current_chunk)
    
    # Post-processing to ensure chunks are not too small or too large
    final_chunks = []
    current_chunk = ""
    
    for chunk in chunks:
        # If chunk is very small, combine with previous unless it's a special section
        if len(chunk) < max_chunk_size / 4 and not any(re.search(pattern, chunk) for pattern in section_patterns):
            if current_chunk and len(current_chunk) + len(chunk) <= max_chunk_size:
                current_chunk += "\n\n" + chunk
            else:
                if current_chunk:
                    final_chunks.append(current_chunk)
                current_chunk = chunk
        else:
            if current_chunk:
                final_chunks.append(current_chunk)
            current_chunk = chunk
    
    # Add the last chunk
    if current_chunk:
        final_chunks.append(current_chunk)
    
    return final_chunks

def chunk_powerpoint_content(text, max_chunk_size=3000, overlap=500):
    """Special chunking function for PowerPoint content that keeps slides together
    
    Args:
        text: The PowerPoint text to chunk
        max_chunk_size: Maximum size of a chunk in characters
        overlap: Number of characters to overlap between chunks
        
    Returns:
        list: List of chunks with multiple slides grouped together
    """
    # If content is small enough, return as is
    if len(text) <= max_chunk_size:
        return [text]
    
    # Split the text by slide markers
    slides = re.split(r'(## SLIDE \d+ ##)', text)
    
    # Pair slide markers with their content
    slide_sections = []
    for i in range(0, len(slides) - 1, 2):
        if i + 1 < len(slides):
            slide_marker = slides[i]
            slide_content = slides[i + 1]
            slide_sections.append(slide_marker + slide_content)
        elif i < len(slides):
            # Handle odd number of splits
            slide_sections.append(slides[i])
    
    # If we didn't get proper slide sections, fallback to simple splitting
    if not slide_sections:
        slide_sections = [text]
    
    # Group slides into chunks
    chunks = []
    current_chunk = ""
    
    for slide in slide_sections:
        # If adding this slide would exceed the max size
        if len(current_chunk) + len(slide) > max_chunk_size and current_chunk:
            chunks.append(current_chunk)
            
            # Get the slides we just added to the chunk for overlap
            slide_markers_in_chunk = re.findall(r'## SLIDE (\d+) ##', current_chunk)
            
            # Start new chunk with presentation info and the last few slides for overlap
            # Aim to include about 2-3 slides for overlap if possible
            if slide_markers_in_chunk:
                # Extract presentation summary if it exists
                presentation_summary = ""
                if "## PRESENTATION SUMMARY ##" in text:
                    summary_match = re.search(r'(## PRESENTATION SUMMARY ##.*?)(?=## SLIDE)', text, re.DOTALL)
                    if summary_match:
                        presentation_summary = summary_match.group(1)
                
                # Find the last 2 slides to include for overlap
                overlap_slide_numbers = slide_markers_in_chunk[-2:] if len(slide_markers_in_chunk) >= 2 else slide_markers_in_chunk
                
                # Start new chunk with presentation summary and overlapping slides
                new_chunk_start = presentation_summary if presentation_summary else ""
                
                # Add overlapping slides
                for slide_num in overlap_slide_numbers:
                    slide_pattern = f"## SLIDE {slide_num} ##.*?(?=## SLIDE|$)"
                    slide_match = re.search(slide_pattern, text, re.DOTALL)
                    if slide_match:
                        new_chunk_start += slide_match.group(0)
                
                # Add the new slide
                current_chunk = new_chunk_start + slide
            else:
                # No slide markers found, use character-based overlap
                overlap_point = max(0, len(current_chunk) - overlap)
                current_chunk = current_chunk[overlap_point:] + slide
        else:
            current_chunk += slide
    
    # Add the last chunk
    if current_chunk:
        chunks.append(current_chunk)
    
    return chunks

def resize_image(image_path, max_dimension=2048):
    """Resize image if it's too large while maintaining aspect ratio"""
    try:
        img = Image.open(image_path)
        
        # Check if resize is needed
        width, height = img.size
        if width <= max_dimension and height <= max_dimension:
            return image_path  # No resize needed
            
        # Calculate new dimensions
        if width > height:
            new_width = max_dimension
            new_height = int(height * (max_dimension / width))
        else:
            new_height = max_dimension
            new_width = int(width * (max_dimension / height))
            
        # Resize and save
        img = img.resize((new_width, new_height), Image.LANCZOS)
        
        # Create new filename for resized image
        filename, ext = os.path.splitext(image_path)
        resized_path = f"{filename}_resized{ext}"
        
        # Save resized image
        img.save(resized_path)
        return resized_path
    except Exception as e:
        print(f"Error resizing image: {str(e)}")
        return image_path  # Return original path if resizing fails 